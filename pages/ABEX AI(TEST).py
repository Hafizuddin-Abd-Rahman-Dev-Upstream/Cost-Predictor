import io
import json
import zipfile
import requests
import numpy as np
import pandas as pd
import streamlit as st
import re

# ML/Stats
from sklearn.impute import KNNImputer, SimpleImputer
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import MinMaxScaler, OneHotEncoder
from sklearn.ensemble import RandomForestRegressor, GradientBoostingRegressor
from sklearn.linear_model import Ridge, Lasso
from sklearn.svm import SVR
from sklearn.tree import DecisionTreeRegressor
from sklearn.pipeline import Pipeline
from sklearn.compose import ColumnTransformer
from sklearn.metrics import mean_squared_error, r2_score
from scipy.stats import linregress

# Viz in-app
import plotly.express as px
import plotly.graph_objects as go

# Viz for PPT
import matplotlib.pyplot as plt

# PPT export
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Excel export helpers
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation

# ---------------------------------------------------------------------------------------
# PAGE CONFIG
# ---------------------------------------------------------------------------------------
st.set_page_config(
    page_title="ABEX AI RT2026",
    page_icon="💠",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------------------------
# THEME TOKENS
# ---------------------------------------------------------------------------------------
PETRONAS = {
    "teal": "#00A19B",
    "teal_dark": "#008C87",
    "purple": "#6C4DD3",
    "white": "#FFFFFF",
    "black": "#0E1116",
    "border": "rgba(0,0,0,0.10)",
}

# ---------------------------------------------------------------------------------------
# SHAREPOINT LINKS (FILL THESE LATER)
# ---------------------------------------------------------------------------------------
SHAREPOINT_LINKS = {
    "Shallow Water": "https://petronas.sharepoint.com/sites/your-site/shallow-water",
    "Deep Water": "https://petronas.sharepoint.com/sites/your-site/deep-water",
    "Onshore": "https://petronas.sharepoint.com/sites/your-site/onshore",
    "Uncon": "https://petronas.sharepoint.com/sites/your-site/uncon",
    "CCS": "https://petronas.sharepoint.com/sites/your-site/ccs",
}

# ---------------------------------------------------------------------------------------
# GLOBAL CSS
# ---------------------------------------------------------------------------------------
st.markdown(
    f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');

html, body {{
  font-family: 'Inter', sans-serif;
}}

[data-testid="stAppViewContainer"] {{
  background: {PETRONAS["white"]};
  color: {PETRONAS["black"]};
  padding-top: 0.5rem;
}}

#MainMenu, footer {{ visibility: hidden; }}

[data-testid="stSidebar"] {{
  background: linear-gradient(180deg, {PETRONAS["teal"]} 0%, {PETRONAS["teal_dark"]} 100%) !important;
  color: #fff !important;
  border-top-right-radius: 16px;
  border-bottom-right-radius: 16px;
  box-shadow: 0 6px 20px rgba(0,0,0,0.15);
}}
[data-testid="stSidebar"] * {{ color: #fff !important; }}

[data-testid="collapsedControl"] {{
  position: fixed !important;
  top: 50% !important;
  left: 10px !important;
  transform: translateY(-50%) !important;
  z-index: 9999 !important;
}}

.petronas-hero {{
  border-radius: 20px;
  padding: 28px 32px;
  margin: 6px 0 18px 0;
  color: #fff;
  background: linear-gradient(135deg, {PETRONAS["teal"]}, {PETRONAS["purple"]}, {PETRONAS["black"]});
  background-size: 200% 200%;
  animation: heroGradient 8s ease-in-out infinite, fadeIn .8s ease-in-out, heroPulse 5s ease-in-out infinite;
  box-shadow: 0 10px 24px rgba(0,0,0,.12);
}}
@keyframes heroGradient {{
  0% {{ background-position: 0% 50%; }}
  50% {{ background-position: 100% 50%; }}
  100% {{ background-position: 0% 50%; }}
}}
@keyframes fadeIn {{
  from {{ opacity: 0; transform: translateY(10px); }}
  to {{ opacity: 1; transform: translateY(0); }}
}}
@keyframes heroPulse {{
  0%   {{ box-shadow: 0 0 16px rgba(0,161,155,0.45); }}
  25%  {{ box-shadow: 0 0 26px rgba(108,77,211,0.55); }}
  50%  {{ box-shadow: 0 0 36px rgba(0,161,155,0.55); }}
  75%  {{ box-shadow: 0 0 26px rgba(108,77,211,0.55); }}
  100% {{ box-shadow: 0 0 16px rgba(0,161,155,0.45); }}
}}
.petronas-hero h1 {{ margin: 0 0 5px; font-weight: 800; letter-spacing: 0.3px; }}
.petronas-hero p {{ margin: 0; opacity: .9; font-weight: 500; }}

.stButton > button, .stDownloadButton > button, .petronas-button {{
  border-radius: 10px;
  padding: .6rem 1.1rem;
  font-weight: 600;
  color: #fff !important;
  border: none;
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  background-size: 200% auto;
  transition: background-position .85s ease, transform .2s ease, box-shadow .25s ease;
  text-decoration: none;
  display: inline-block;
}}
.stButton > button:hover, .stDownloadButton > button:hover, .petronas-button:hover {{
  background-position: right center;
  transform: translateY(-1px);
  box-shadow: 0 6px 16px rgba(0,0,0,0.18);
}}

.stTabs [role="tablist"] {{
  display: flex;
  gap: 8px;
  border-bottom: none;
  padding-bottom: 6px;
}}
.stTabs [role="tab"] {{
  background: #fff;
  color: {PETRONAS["black"]};
  border-radius: 8px;
  padding: 10px 18px;
  border: 1px solid {PETRONAS["border"]};
  font-weight: 600;
  transition: all .3s ease;
  position: relative;
}}
.stTabs [role="tab"]:hover {{
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  color: #fff;
}}
.stTabs [role="tab"][aria-selected="true"] {{
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  color: #fff;
  border-color: transparent;
  box-shadow: 0 4px 16px rgba(0,0,0,0.15);
}}
.stTabs [role="tab"][aria-selected="true"]::after {{
  content: "";
  position: absolute;
  left: 10%;
  bottom: -3px;
  width: 80%;
  height: 3px;
  background: linear-gradient(90deg, {PETRONAS["teal"]}, {PETRONAS["purple"]}, {PETRONAS["teal"]});
  background-size: 200% 100%;
  border-radius: 2px;
  animation: glowSlide 2.5s linear infinite;
}}
@keyframes glowSlide {{
  0% {{ background-position: 0% 50%; }}
  50% {{ background-position: 100% 50%; }}
  100% {{ background-position: 0% 50%; }}
}}
</style>
""",
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------------------------
# HERO HEADER
# ---------------------------------------------------------------------------------------
st.markdown(
    """
<div class="petronas-hero">
  <h1>ABEX AI RT2026</h1>
  <p>Data-driven ABEX prediction</p>
</div>
""",
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------------------------
# AUTH
# ---------------------------------------------------------------------------------------
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False

APPROVED_EMAILS = st.secrets.get("emails", [])
correct_password = st.secrets.get("password", None)

if not st.session_state.authenticated:
    with st.form("login_form"):
        st.markdown("#### 🔐 Access Required", unsafe_allow_html=True)
        email = st.text_input("Email Address")
        password = st.text_input("Access Password", type="password")
        submitted = st.form_submit_button("Login")

        if submitted:
            if email in APPROVED_EMAILS and password == correct_password:
                st.session_state.authenticated = True
                st.success("✅ Access granted.")
                st.rerun()
            else:
                st.error("❌ Invalid email or password.")
    st.stop()

# ---------------------------------------------------------------------------------------
# SESSION STATE
# ---------------------------------------------------------------------------------------
if "datasets" not in st.session_state:
    st.session_state.datasets = {}
if "predictions" not in st.session_state:
    st.session_state.predictions = {}
if "processed_excel_files" not in st.session_state:
    st.session_state.processed_excel_files = set()
if "_last_metrics" not in st.session_state:
    st.session_state._last_metrics = None
if "projects" not in st.session_state:
    st.session_state.projects = {}
if "component_labels" not in st.session_state:
    st.session_state.component_labels = {}
if "best_model_name_per_dataset" not in st.session_state:
    st.session_state.best_model_name_per_dataset = {}

# Floater specific session state
if "floater_dataset" not in st.session_state:
    st.session_state.floater_dataset = None
if "floater_model" not in st.session_state:
    st.session_state.floater_model = None
if "floater_metrics" not in st.session_state:
    st.session_state.floater_metrics = None
if "floater_predictions" not in st.session_state:
    st.session_state.floater_predictions = []
if "floater_feature_columns" not in st.session_state:
    st.session_state.floater_feature_columns = []

# IMPORTANT: uploader reset nonce (so "clear uploaded files" truly clears the widget)
if "uploader_nonce" not in st.session_state:
    st.session_state.uploader_nonce = 0

# ---------------------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------------------
def toast(msg, icon="✅"):
    try:
        st.toast(f"{icon} {msg}")
    except Exception:
        st.success(msg if icon == "✅" else msg)


def format_with_commas(num):
    try:
        return f"{float(num):,.2f}"
    except Exception:
        return str(num)


# =========================
# FIXED: currency detection
# =========================
def get_currency_symbol(df: pd.DataFrame) -> str:
    """
    Detect currency from the target cost column.

    Your rule:
    - Currency column is usually the LAST column (e.g. "Total Cost (Mil USD)").
    - If there are multiple "Total Cost ..." columns, choose the LAST "Total Cost" column.
    - Fallback: choose the last non-junk column.
    """
    if df is None or df.empty:
        return ""

    def is_junk(colname: str) -> bool:
        h = str(colname).strip().upper()
        return (not h) or h.startswith("UNNAMED") or h in {"INDEX", "IDX"}

    # 1) Prefer LAST "Total Cost" column
    cost_cols = []
    for c in df.columns:
        h = str(c).strip().upper()
        if not is_junk(c) and ("TOTAL" in h and "COST" in h):
            cost_cols.append(c)

    if cost_cols:
        preferred = cost_cols[-1]  # <-- LAST Total Cost column
    else:
        # 2) Fallback: last non-junk column
        preferred = None
        for c in reversed(df.columns):
            if not is_junk(c):
                preferred = c
                break

    if preferred is None:
        return ""

    header = str(preferred).strip().upper()

    # symbols
    if "€" in header:
        return "€"
    if "£" in header:
        return "£"
    if "$" in header:
        return "USD"

    # codes (strict word boundaries)
    if re.search(r"\bUSD\b", header):
        return "USD"
    if re.search(r"\b(MYR|RM)\b", header):
        return "RM"

    return ""


def normalize_to_100(d: dict):
    total = sum(float(v) for v in d.values())
    if total <= 0:
        return d, total
    out = {k: float(v) * 100.0 / total for k, v in d.items()}
    keys = list(out.keys())
    rounded = {k: round(out[k], 2) for k in keys}
    diff = 100.0 - sum(rounded.values())
    if keys:
        rounded[keys[-1]] = round(rounded[keys[-1]] + diff, 2)
    return rounded, total


def cost_breakdown(
    base_pred: float,
    eprr: dict,
    sst_pct: float,
    owners_pct: float,
    cont_pct: float,
    esc_pct: float,
):
    owners_cost = round(float(base_pred) * (owners_pct / 100.0), 2)
    sst_cost = round(float(base_pred) * (sst_pct / 100.0), 2)
    contingency_cost = round((float(base_pred) + owners_cost) * (cont_pct / 100.0), 2)
    escalation_cost = round((float(base_pred) + owners_cost) * (esc_pct / 100.0), 2)

    eprr_costs = {k: round(float(base_pred) * (float(v) / 100.0), 2) for k, v in (eprr or {}).items()}
    grand_total = round(float(base_pred) + owners_cost + contingency_cost + escalation_cost + sst_cost, 2)

    return owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total


def project_components_df(proj):
    comps = proj.get("components", [])
    rows = []
    for c in comps:
        rows.append(
            {
                "Component": c["component_type"],
                "Dataset": c["dataset"],
                "Base ABEX": float(c["prediction"]),
                "Owner's Cost": float(c["breakdown"]["owners_cost"]),
                "Contingency": float(c["breakdown"]["contingency_cost"]),
                "Escalation": float(c["breakdown"]["escalation_cost"]),
                "SST": float(c["breakdown"]["sst_cost"]),
                "Grand Total": float(c["breakdown"]["grand_total"]),
            }
        )
    return pd.DataFrame(rows)


def create_project_excel_report_ABEX(project_name, proj, currency=""):
    output = io.BytesIO()
    comps_df = project_components_df(proj)

    if comps_df.empty:
        with pd.ExcelWriter(output, engine="openpyxl") as writer:
            pd.DataFrame({"Info": [f"No components for project {project_name}"]}).to_excel(
                writer, sheet_name="Summary", index=False
            )
        output.seek(0)
        return output

    total_ABEX = comps_df["Base ABEX"].sum()
    total_grand = comps_df["Grand Total"].sum()

    summary_df = comps_df.copy()
    summary_df.loc[len(summary_df)] = {
        "Component": "TOTAL",
        "Dataset": "",
        "Base ABEX": total_ABEX,
        "Owner's Cost": comps_df["Owner's Cost"].sum(),
        "Contingency": comps_df["Contingency"].sum(),
        "Escalation": comps_df["Escalation"].sum(),
        "SST": comps_df["SST"].sum(),
        "Grand Total": total_grand,
    }

    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df.to_excel(writer, sheet_name="Summary", index=False)
        ws = writer.sheets["Summary"]

        max_row = ws.max_row
        max_col = ws.max_column

        for col_idx in range(3, max_col + 1):
            col_letter = get_column_letter(col_idx)
            ws.conditional_formatting.add(
                f"{col_letter}2:{col_letter}{max_row-1}",
                ColorScaleRule(
                    start_type="percentile",
                    start_value=10,
                    start_color="FFE0F7FA",
                    mid_type="percentile",
                    mid_value=50,
                    mid_color="FF80DEEA",
                    end_type="percentile",
                    end_value=90,
                    end_color="FF00838F",
                ),
            )

        chart = BarChart()
        chart.title = "Grand Total by Component"
        data = Reference(ws, min_col=8, max_col=8, min_row=1, max_row=max_row - 1)
        cats = Reference(ws, min_col=1, min_row=2, max_row=max_row - 1)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.y_axis.title = f"Cost ({currency})".strip()
        chart.x_axis.title = "Component"
        chart.height = 10
        chart.width = 18
        ws.add_chart(chart, "J2")

        line = LineChart()
        line.title = "Base ABEX Trend"
        data_ABEX = Reference(ws, min_col=3, max_col=3, min_row=1, max_row=max_row - 1)
        line.add_data(data_ABEX, titles_from_data=True)
        line.set_categories(cats)
        line.y_axis.title = f"Base ABEX ({currency})".strip()
        line.height = 10
        line.width = 18
        ws.add_chart(line, "J20")

        comps_df.to_excel(writer, sheet_name="Components Detail", index=False)

    output.seek(0)
    return output


def create_project_pptx_report_ABEX(project_name, proj, currency=""):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    layout_title_only = prs.slide_layouts[5]
    layout_title_content = prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout_title_only)
    title = slide.shapes.title
    title.text = f"ABEX Project Report\n{project_name}"
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 161, 155)

    comps_df = project_components_df(proj)
    comps = proj.get("components", [])
    total_ABEX = comps_df["Base ABEX"].sum() if not comps_df.empty else 0.0
    total_grand = comps_df["Grand Total"].sum() if not comps_df.empty else 0.0

    slide = prs.slides.add_slide(layout_title_content)
    slide.shapes.title.text = "Executive Summary"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    lines = [
        f"Project: {project_name}",
        f"Total Components: {len(comps)}",
        f"Total Base ABEX: {currency} {total_ABEX:,.2f}",
        f"Total Grand Total: {currency} {total_grand:,.2f}",
        "",
        "Components:",
    ]
    for c in comps:
        lines.append(f"• {c['component_type']}: {currency} {c['breakdown']['grand_total']:,.2f}")

    for i, line in enumerate(lines):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.text = line
        para.font.size = Pt(16)

    if not comps_df.empty:
        fig, ax = plt.subplots(figsize=(7, 4))
        ax.bar(comps_df["Component"], comps_df["Grand Total"])
        ax.set_title("Grand Total by Component")
        ax.set_ylabel(f"Cost ({currency})".strip())
        ax.tick_params(axis="x", rotation=25)
        ax.grid(axis="y", linestyle="--", alpha=0.4)
        fig.tight_layout()

        img_stream = io.BytesIO()
        fig.savefig(img_stream, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        img_stream.seek(0)

        slide = prs.slides.add_slide(layout_title_only)
        slide.shapes.title.text = "Grand Total by Component"
        slide.shapes.add_picture(img_stream, Inches(0.7), Inches(1.5), width=Inches(8.6))

        fig2, ax2 = plt.subplots(figsize=(7, 4))
        labels = comps_df["Component"]
        base = comps_df["Base ABEX"]
        owners = comps_df["Owner's Cost"]
        cont = comps_df["Contingency"]
        esc = comps_df["Escalation"]
        sst = comps_df["SST"]

        bottom = np.zeros(len(labels))
        for vals, lab in [
            (base, "Base ABEX"),
            (owners, "Owner"),
            (cont, "Contingency"),
            (esc, "Escalation"),
            (sst, "SST"),
        ]:
            ax2.bar(labels, vals, bottom=bottom, label=lab)
            bottom += np.array(vals)

        ax2.set_title("Cost Composition by Component")
        ax2.set_ylabel(f"Cost ({currency})".strip())
        ax2.tick_params(axis="x", rotation=25)
        ax2.grid(axis="y", linestyle="--", alpha=0.4)
        ax2.legend(fontsize=8, ncol=3)
        fig2.tight_layout()

        img_stream2 = io.BytesIO()
        fig2.savefig(img_stream2, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig2)
        img_stream2.seek(0)

        slide2 = prs.slides.add_slide(layout_title_only)
        slide2.shapes.title.text = "Cost Composition by Component"
        slide2.shapes.add_picture(img_stream2, Inches(0.7), Inches(1.5), width=Inches(8.6))

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def create_comparison_excel_report_ABEX(projects_dict, currency=""):
    output = io.BytesIO()

    summary_rows = []
    for name, proj in projects_dict.items():
        dfc = project_components_df(proj)
        ABEX = dfc["Base ABEX"].sum() if not dfc.empty else 0.0
        owners = dfc["Owner's Cost"].sum() if not dfc.empty else 0.0
        cont = dfc["Contingency"].sum() if not dfc.empty else 0.0
        esc = dfc["Escalation"].sum() if not dfc.empty else 0.0
        sst = dfc["SST"].sum() if not dfc.empty else 0.0
        grand = dfc["Grand Total"].sum() if not dfc.empty else 0.0
        summary_rows.append(
            {
                "Project": name,
                "Components": len(proj.get("components", [])),
                "ABEX Sum": ABEX,
                "Owner": owners,
                "Contingency": cont,
                "Escalation": esc,
                "SST": sst,
                "Grand Total": grand,
            }
        )

    summary_df = pd.DataFrame(summary_rows)

    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df.to_excel(writer, sheet_name="Projects Summary", index=False)
        ws = writer.sheets["Projects Summary"]

        max_row = ws.max_row
        max_col = ws.max_column

        for col_idx in range(3, max_col + 1):
            col_letter = get_column_letter(col_idx)
            ws.conditional_formatting.add(
                f"{col_letter}2:{col_letter}{max_row}",
                ColorScaleRule(
                    start_type="percentile",
                    start_value=10,
                    start_color="FFE3F2FD",
                    mid_type="percentile",
                    mid_value=50,
                    mid_color="FF90CAF9",
                    end_type="percentile",
                    end_value=90,
                    end_color="FF1565C0",
                ),
            )

        chart = BarChart()
        chart.title = "Grand Total by Project"
        data = Reference(ws, min_col=8, max_col=8, min_row=1, max_row=max_row)
        cats = Reference(ws, min_col=1, min_row=2, max_row=max_row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.y_axis.title = f"Cost ({currency})".strip()
        chart.x_axis.title = "Project"
        chart.height = 10
        chart.width = 18
        ws.add_chart(chart, "J2")

        for name, proj in projects_dict.items():
            dfc = project_components_df(proj)
            if dfc.empty:
                continue
            sheet_name = name[:31]
            dfc.to_excel(writer, sheet_name=sheet_name, index=False)

    output.seek(0)
    return output


def create_comparison_pptx_report_ABEX(projects_dict, currency=""):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    layout_title_only = prs.slide_layouts[5]

    slide = prs.slides.add_slide(layout_title_only)
    title = slide.shapes.title
    title.text = "ABEX Project Comparison"
    p = title.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 161, 155)

    rows = []
    for name, proj in projects_dict.items():
        dfc = project_components_df(proj)
        ABEX = dfc["Base ABEX"].sum() if not dfc.empty else 0.0
        owners = dfc["Owner's Cost"].sum() if not dfc.empty else 0.0
        cont = dfc["Contingency"].sum() if not dfc.empty else 0.0
        esc = dfc["Escalation"].sum() if not dfc.empty else 0.0
        sst = dfc["SST"].sum() if not dfc.empty else 0.0
        grand = dfc["Grand Total"].sum() if not dfc.empty else 0.0
        rows.append(
            {
                "Project": name,
                "ABEX Sum": ABEX,
                "Owner": owners,
                "Contingency": cont,
                "Escalation": esc,
                "SST": sst,
                "Grand Total": grand,
            }
        )
    df_proj = pd.DataFrame(rows)

    if not df_proj.empty:
        fig, ax = plt.subplots(figsize=(7, 4))
        ax.bar(df_proj["Project"], df_proj["Grand Total"])
        ax.set_title("Grand Total by Project")
        ax.set_ylabel(f"Cost ({currency})".strip())
        ax.tick_params(axis="x", rotation=25)
        ax.grid(axis="y", linestyle="--", alpha=0.4)
        fig.tight_layout()

        img_stream = io.BytesIO()
        fig.savefig(img_stream, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        img_stream.seek(0)

        slide = prs.slides.add_slide(layout_title_only)
        slide.shapes.title.text = "Grand Total by Project"
        slide.shapes.add_picture(img_stream, Inches(0.7), Inches(1.5), width=Inches(8.6))

        fig2, ax2 = plt.subplots(figsize=(7, 4))
        labels = df_proj["Project"]
        base = df_proj["ABEX Sum"]
        owners = df_proj["Owner"]
        cont = df_proj["Contingency"]
        esc = df_proj["Escalation"]
        sst = df_proj["SST"]

        bottom = np.zeros(len(labels))
        for vals, lab in [
            (base, "Base ABEX"),
            (owners, "Owner"),
            (cont, "Contingency"),
            (esc, "Escalation"),
            (sst, "SST"),
        ]:
            ax2.bar(labels, vals, bottom=bottom, label=lab)
            bottom += np.array(vals)

        ax2.set_title("Cost Composition by Project")
        ax2.set_ylabel(f"Cost ({currency})".strip())
        ax2.tick_params(axis="x", rotation=25)
        ax2.grid(axis="y", linestyle="--", alpha=0.4)
        ax2.legend(fontsize=8, ncol=3)
        fig2.tight_layout()

        img_stream2 = io.BytesIO()
        fig2.savefig(img_stream2, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig2)
        img_stream2.seek(0)

        slide2 = prs.slides.add_slide(layout_title_only)
        slide2.shapes.title.text = "Cost Composition by Project"
        slide2.shapes.add_picture(img_stream2, Inches(0.7), Inches(1.5), width=Inches(8.6))

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ---------------------------------------------------------------------------------------
# DATA / MODEL HELPERS
# ---------------------------------------------------------------------------------------
GITHUB_USER = "Hafizuddin-Abd-Rahman-Dev-Upstream"
REPO_NAME = "Cost-Predictor"
BRANCH = "main"
DATA_FOLDER = "pages/data_ABEX"

MODEL_CANDIDATES = {
    "RandomForest": lambda rs=42: RandomForestRegressor(random_state=rs),
    "GradientBoosting": lambda rs=42: GradientBoostingRegressor(random_state=rs),
    "Ridge": lambda rs=42: Ridge(),
    "Lasso": lambda rs=42: Lasso(),
    "SVR": lambda rs=42: SVR(),
    "DecisionTree": lambda rs=42: DecisionTreeRegressor(random_state=rs),
}


@st.cache_data(ttl=600)
def list_csvs_from_manifest(folder_path: str):
    manifest_url = f"https://raw.githubusercontent.com/{GITHUB_USER}/{REPO_NAME}/{BRANCH}/{folder_path}/files.json"
    try:
        res = requests.get(manifest_url, timeout=10)
        res.raise_for_status()
        return res.json()
    except Exception as e:
        st.error(f"Failed to load CSV manifest: {e}")
        return []


# ---------------------------------------------------------------------------------------
# NAV ROW — FIVE SHAREPOINT BUTTONS
# ---------------------------------------------------------------------------------------
nav_labels = ["SHALLOW WATER", "DEEP WATER", "ONSHORE", "UNCON", "CCS"]
nav_cols = st.columns(len(nav_labels))
for col, label in zip(nav_cols, nav_labels):
    with col:
        url = SHAREPOINT_LINKS.get(label.title(), "#")
        st.markdown(
            f'''
            <a href="{url}" target="_blank" rel="noopener"
               class="petronas-button"
               style="width:100%; text-align:center; display:inline-block;">
               {label}
            </a>
            ''',
            unsafe_allow_html=True,
        )

# ---------------------------------------------------------------------------------------
# TOP-LEVEL TABS
# ---------------------------------------------------------------------------------------
tab_data, tab_pb, tab_compare = st.tabs(["📊 Data", "🏗️ Project Builder", "🔀 Compare Projects"])


# =======================================================================================
# DATA TAB
# =======================================================================================
with tab_data:
    st.markdown('<h3 style="margin-top:0;color:#000;">📁 Data</h3>', unsafe_allow_html=True)

    st.markdown('<h4 style="margin:0;color:#000;">Data Sources</h4><p></p>', unsafe_allow_html=True)
    c1, c2 = st.columns([1.2, 1])
    with c1:
        data_source = st.radio("Choose data source", ["Upload CSV", "Load from Server"], horizontal=True)

    with c2:
        st.caption("Enterprise Storage (SharePoint)")
        data_link = (
            "https://petronas.sharepoint.com/sites/ecm_ups_coe/confidential/"
            "DFE%20Cost%20Engineering/Forms/AllItems.aspx?"
            "id=%2Fsites%2Fecm%5Fups%5Fcoe%2Fconfidential%2FDFE%20Cost%20Engineering"
            "%2F2%2ETemplate%20Tools%2FCost%20Predictor%2FDatabase%2FABEX%20%2D%20RT%20Q1%202025"
            "&viewid=25092e6d%2D373d%2D41fe%2D8f6f%2D486cd8cdd5b8"
        )
        st.markdown(
            f'<a href="{data_link}" target="_blank" rel="noopener" class="petronas-button">Open Enterprise Storage</a>',
            unsafe_allow_html=True,
        )

    uploaded_files = []
    if data_source == "Upload CSV":
        uploaded_files = st.file_uploader(
            "Upload CSV files (max 200MB)",
            type="csv",
            accept_multiple_files=True,
            key=f"csv_uploader_{st.session_state.uploader_nonce}",
        )
    else:
        github_csvs = list_csvs_from_manifest(DATA_FOLDER)
        if github_csvs:
            selected_file = st.selectbox("Choose CSV from GitHub", github_csvs)
            if st.button("Load selected CSV"):
                raw_url = f"https://raw.githubusercontent.com/{GITHUB_USER}/{REPO_NAME}/{BRANCH}/{DATA_FOLDER}/{selected_file}"
                try:
                    df = pd.read_csv(raw_url)
                    st.session_state.datasets[selected_file] = df
                    st.session_state.predictions.setdefault(selected_file, [])
                    toast(f"Loaded from GitHub: {selected_file}")
                    st.rerun()
                except Exception as e:
                    st.error(f"Error loading CSV: {e}")
        else:
            st.info("No CSV files found in GitHub folder.")

    if uploaded_files:
        for up in uploaded_files:
            if up.name not in st.session_state.datasets:
                try:
                    df = pd.read_csv(up)
                    st.session_state.datasets[up.name] = df
                    st.session_state.predictions.setdefault(up.name, [])
                except Exception as e:
                    st.error(f"Failed to read {up.name}: {e}")
        toast("Dataset(s) added.")

    st.divider()

    cA, cB, cC, cD = st.columns([1, 1, 1, 2])
    with cA:
        if st.button("🧹 Clear all predictions"):
            st.session_state.predictions = {k: [] for k in st.session_state.predictions.keys()}
            toast("All predictions cleared.", "🧹")
            st.rerun()

    with cB:
        if st.button("🧺 Clear processed files history"):
            st.session_state.processed_excel_files = set()
            toast("Processed files history cleared.", "🧺")
            st.rerun()

    with cC:
        if st.button("🔁 Refresh server manifest"):
            list_csvs_from_manifest.clear()
            toast("Server manifest refreshed.", "🔁")
            st.rerun()

    with cD:
        if st.button("🗂️ Clear all uploaded / loaded files (keep projects)"):
            st.session_state.datasets = {}
            st.session_state.predictions = {}
            st.session_state.processed_excel_files = set()
            st.session_state._last_metrics = None
            st.session_state.best_model_name_per_dataset = {}

            for k in ["ds_model", "ds_viz", "ds_pred", "ds_results"]:
                if k in st.session_state:
                    del st.session_state[k]

            st.session_state.uploader_nonce += 1
            toast("All datasets cleared. Projects preserved.", "🗂️")
            st.rerun()

    st.divider()

    if st.session_state.datasets:
        ds_name_data = st.selectbox("Active dataset", list(st.session_state.datasets.keys()))
        df_active = st.session_state.datasets[ds_name_data]
        currency_active = get_currency_symbol(df_active)
        colA, colB, colC = st.columns([1, 1, 1])
        with colA:
            st.metric("Rows", f"{df_active.shape[0]:,}")
        with colB:
            st.metric("Columns", f"{df_active.shape[1]:,}")
        with colC:
            st.metric("Currency", f"{currency_active or '—'}")
        with st.expander("Preview (first 10 rows)", expanded=False):
            st.dataframe(df_active.head(10), use_container_width=True)
    else:
        st.info("Upload or load a dataset to proceed.")

    # ========================= MODEL TRAINING =================================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">⚙️ Model</h3>', unsafe_allow_html=True)

    if not st.session_state.datasets:
        st.info("No dataset. Use the Data section above to upload or load.")
    else:
        ds_name_model = st.selectbox("Dataset for model training", list(st.session_state.datasets.keys()), key="ds_model")
        df_model = st.session_state.datasets[ds_name_model]

        with st.spinner("Preparing data..."):
            # Separate features and target (assume last column is target)
            X_raw = df_model.iloc[:, :-1]
            y_raw = df_model.iloc[:, -1]

            # Check if target is numeric
            if not pd.api.types.is_numeric_dtype(y_raw):
                st.error("Target column (last column) must be numeric. Please check your dataset.")
                st.stop()

            # Identify column types
            numeric_cols = X_raw.select_dtypes(include=np.number).columns.tolist()
            categorical_cols = X_raw.select_dtypes(include='object').columns.tolist()

            # Impute missing values only on numeric columns
            if numeric_cols:
                imputer_num = KNNImputer(n_neighbors=5)
                X_numeric_imputed = imputer_num.fit_transform(X_raw[numeric_cols])
                X_numeric_imputed = pd.DataFrame(X_numeric_imputed, columns=numeric_cols, index=X_raw.index)
            else:
                X_numeric_imputed = pd.DataFrame()

            # Keep categorical columns as they are (could fill mode if desired)
            X_cat = X_raw[categorical_cols] if categorical_cols else pd.DataFrame()

            # Combine back in original order
            X_processed = pd.concat([X_numeric_imputed, X_cat], axis=1)
            X_processed = X_processed[X_raw.columns]  # restore original column order
            y_processed = y_raw

        st.markdown('<h4 style="margin:0;color:#000;">Train & Evaluate</h4><p>Step 2</p>', unsafe_allow_html=True)
        m1, m2 = st.columns([1, 3])
        with m1:
            test_size = st.slider("Test size", 0.1, 0.5, 0.2, 0.05)
            run_train = st.button("Run training")
        with m2:
            st.caption("Automatic best-model selection over 6 regressors (with ColumnTransformer for mixed types).")

        if run_train:
            with st.spinner("Training model..."):
                # Split data
                X_train, X_test, y_train, y_test = train_test_split(
                    X_processed, y_processed, test_size=test_size, random_state=42
                )

                # Define preprocessor for numeric and categorical columns
                numeric_transformer = Pipeline(steps=[
                    ('imputer', SimpleImputer(strategy='median')),  # just in case
                    ('scaler', MinMaxScaler())
                ])
                categorical_transformer = Pipeline(steps=[
                    ('imputer', SimpleImputer(strategy='constant', fill_value='missing')),
                    ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
                ])

                # Update column lists based on the actual columns present after imputation (should be same)
                num_cols_train = X_train.select_dtypes(include=np.number).columns.tolist()
                cat_cols_train = X_train.select_dtypes(include='object').columns.tolist()

                preprocessor = ColumnTransformer(
                    transformers=[
                        ('num', numeric_transformer, num_cols_train),
                        ('cat', categorical_transformer, cat_cols_train)
                    ]
                )

                # Evaluate each model
                rows = []
                best_name = None
                best_r2 = -np.inf
                best_rmse = None

                for name, ctor in MODEL_CANDIDATES.items():
                    try:
                        base_model = ctor(42)
                    except TypeError:
                        base_model = ctor()

                    pipeline = Pipeline(steps=[
                        ('preprocessor', preprocessor),
                        ('model', base_model)
                    ])

                    pipeline.fit(X_train, y_train)
                    y_pred = pipeline.predict(X_test)
                    rmse = np.sqrt(mean_squared_error(y_test, y_pred))
                    r2 = r2_score(y_test, y_pred)

                    rows.append({"model": name, "rmse": rmse, "r2": r2})
                    if r2 > best_r2:
                        best_r2 = r2
                        best_rmse = rmse
                        best_name = name

                # Store best model info
                metrics = {
                    "best_model": best_name,
                    "rmse": best_rmse,
                    "r2": best_r2,
                    "models": sorted(rows, key=lambda d: d["r2"], reverse=True)
                }
                st.session_state._last_metrics = metrics
                st.session_state.best_model_name_per_dataset[ds_name_model] = best_name

                # Also store the column types for later use in prediction
                st.session_state[f"coltypes_{ds_name_model}"] = {
                    'numeric': num_cols_train,
                    'categorical': cat_cols_train
                }

                # Display results
                m3, m4 = st.columns(2)
                with m3:
                    st.metric("RMSE (best)", f"{best_rmse:,.2f}")
                with m4:
                    st.metric("R² (best)", f"{best_r2:.3f}")

                toast("Training complete.")
                st.caption(f"Best model selected: **{best_name}**")

                # Show comparison table
                df_models = pd.DataFrame(rows).set_index("model")
                st.markdown("##### Model comparison (6-model pool)")
                styled = (
                    df_models.style.format({"rmse": "{:,.2f}", "r2": "{:.3f}"})
                    .background_gradient(subset=["r2"], cmap="YlGn")
                    .background_gradient(subset=["rmse"], cmap="OrRd_r")
                )
                st.dataframe(styled, use_container_width=True)

    # ========================= VISUALIZATION ==================================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">📈 Visualization</h3>', unsafe_allow_html=True)

    if not st.session_state.datasets:
        st.info("No dataset. Use the Data section above to upload or load.")
    else:
        ds_name_viz = st.selectbox("Dataset for visualization", list(st.session_state.datasets.keys()), key="ds_viz")
        df_viz = st.session_state.datasets[ds_name_viz]

        # Prepare data (same as training: impute numeric, keep categorical)
        X_raw_viz = df_viz.iloc[:, :-1]
        y_viz = df_viz.iloc[:, -1]
        numeric_cols_viz = X_raw_viz.select_dtypes(include=np.number).columns.tolist()
        categorical_cols_viz = X_raw_viz.select_dtypes(include='object').columns.tolist()

        if numeric_cols_viz:
            imputer_num = KNNImputer(n_neighbors=5)
            X_num_imp = imputer_num.fit_transform(X_raw_viz[numeric_cols_viz])
            X_num_imp = pd.DataFrame(X_num_imp, columns=numeric_cols_viz, index=X_raw_viz.index)
        else:
            X_num_imp = pd.DataFrame()

        X_cat_viz = X_raw_viz[categorical_cols_viz] if categorical_cols_viz else pd.DataFrame()
        X_processed_viz = pd.concat([X_num_imp, X_cat_viz], axis=1)
        X_processed_viz = X_processed_viz[X_raw_viz.columns]

        st.markdown('<h4 style="margin:0;color:#000;">Correlation Matrix</h4><p>Exploration</p>', unsafe_allow_html=True)
        if numeric_cols_viz:
            corr = X_processed_viz[numeric_cols_viz].corr()
            fig_corr = px.imshow(corr, text_auto=".2f", aspect="auto", color_continuous_scale="RdBu_r", zmin=-1, zmax=1)
            fig_corr.update_layout(margin=dict(l=0, r=0, t=10, b=0))
            st.plotly_chart(fig_corr, use_container_width=True)
        else:
            st.info("No numeric columns for correlation matrix.")

        st.markdown('<h4 style="margin:0;color:#000;">Feature Importance</h4><p>Model</p>', unsafe_allow_html=True)
        if X_processed_viz.shape[1] > 0:
            # Create a simple preprocessor and model just for importance
            num_cols = X_processed_viz.select_dtypes(include=np.number).columns.tolist()
            cat_cols = X_processed_viz.select_dtypes(include='object').columns.tolist()

            if cat_cols:
                # Use ColumnTransformer with one-hot
                preprocessor = ColumnTransformer([
                    ('num', Pipeline([('scaler', MinMaxScaler())]), num_cols),
                    ('cat', OneHotEncoder(handle_unknown='ignore', sparse_output=False), cat_cols)
                ])
                model = RandomForestRegressor(random_state=42)
                pipeline = Pipeline([('pre', preprocessor), ('model', model)])
                pipeline.fit(X_processed_viz, y_viz)
                # Get feature names after one-hot
                feature_names = []
                feature_names.extend(num_cols)
                ohe = pipeline.named_steps['pre'].named_transformers_['cat']
                feature_names.extend(ohe.get_feature_names_out(cat_cols))
                importances = pipeline.named_steps['model'].feature_importances_
                # Truncate to actual length (some features might be dropped if constant)
                if len(importances) < len(feature_names):
                    feature_names = feature_names[:len(importances)]
                fi = pd.DataFrame({"feature": feature_names, "importance": importances}).sort_values("importance", ascending=True)
            else:
                # All numeric
                scaler = MinMaxScaler()
                X_scaled = scaler.fit_transform(X_processed_viz)
                model = RandomForestRegressor(random_state=42).fit(X_scaled, y_viz)
                importances = model.feature_importances_
                fi = pd.DataFrame({"feature": X_processed_viz.columns, "importance": importances}).sort_values("importance", ascending=True)

            fig_fi = go.Figure(go.Bar(x=fi["importance"], y=fi["feature"], orientation="h"))
            fig_fi.update_layout(xaxis_title="Importance", yaxis_title="Feature", margin=dict(l=0, r=0, t=10, b=0))
            st.plotly_chart(fig_fi, use_container_width=True)
        else:
            st.warning("No features available for importance plot.")

        st.markdown('<h4 style="margin:0;color:#000;">Cost Curve</h4><p>Trend</p>', unsafe_allow_html=True)
        if numeric_cols_viz:
            feat = st.selectbox("Select feature for cost curve", numeric_cols_viz)
            x_vals = X_processed_viz[feat].values
            y_vals = y_viz.values
            mask = (~np.isnan(x_vals)) & (~np.isnan(y_vals))
            scatter_df = pd.DataFrame({feat: x_vals[mask], y_viz.name: y_vals[mask]})
            fig_cc = px.scatter(scatter_df, x=feat, y=y_viz.name, opacity=0.65)

            if mask.sum() >= 2 and np.unique(x_vals[mask]).size >= 2:
                xv = scatter_df[feat].to_numpy(dtype=float)
                yv = scatter_df[y_viz.name].to_numpy(dtype=float)
                slope, intercept, r_value, p_value, std_err = linregress(xv, yv)
                x_line = np.linspace(xv.min(), xv.max(), 100)
                y_line = slope * x_line + intercept
                fig_cc.add_trace(
                    go.Scatter(
                        x=x_line,
                        y=y_line,
                        mode="lines",
                        name=f"Fit: y={slope:.2f}x+{intercept:.2f} (R²={r_value**2:.3f})",
                    )
                )
            else:
                st.warning("Not enough valid/variable data to compute regression.")
            fig_cc.update_layout(margin=dict(l=0, r=0, t=10, b=0))
            st.plotly_chart(fig_cc, use_container_width=True)
        else:
            st.info("No numeric features for cost curve.")

    # ========================= PREDICT =======================================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">🎯 Predict</h3>', unsafe_allow_html=True)

    if not st.session_state.datasets:
        st.info("No dataset. Use the Data section above to upload or load.")
    else:
        ds_name_pred = st.selectbox("Dataset for prediction", list(st.session_state.datasets.keys()), key="ds_pred")
        df_pred = st.session_state.datasets[ds_name_pred]
        currency_pred = get_currency_symbol(df_pred)

        # Prepare features and target (same as training)
        X_raw_pred = df_pred.iloc[:, :-1]
        y_raw_pred = df_pred.iloc[:, -1]
        numeric_cols_pred = X_raw_pred.select_dtypes(include=np.number).columns.tolist()
        categorical_cols_pred = X_raw_pred.select_dtypes(include='object').columns.tolist()

        if numeric_cols_pred:
            imputer_num = KNNImputer(n_neighbors=5)
            X_num_imp = imputer_num.fit_transform(X_raw_pred[numeric_cols_pred])
            X_num_imp = pd.DataFrame(X_num_imp, columns=numeric_cols_pred, index=X_raw_pred.index)
        else:
            X_num_imp = pd.DataFrame()

        X_cat_pred = X_raw_pred[categorical_cols_pred] if categorical_cols_pred else pd.DataFrame()
        X_processed_pred = pd.concat([X_num_imp, X_cat_pred], axis=1)
        X_processed_pred = X_processed_pred[X_raw_pred.columns]
        y_processed_pred = y_raw_pred
        target_column_pred = y_processed_pred.name

        st.markdown('<h4 style="margin:0;color:#000;">Configuration (EPRR • Financial)</h4><p>Step 3</p>', unsafe_allow_html=True)

        c1, c2 = st.columns([1, 1])
        with c1:
            st.markdown("**EPRR Breakdown (%)** (use +/-)")
            eng = st.number_input("Engineering (%)", min_value=0.0, max_value=100.0, value=12.0, step=1.0)
            prep = st.number_input("Preparation (%)", min_value=0.0, max_value=100.0, value=7.0, step=1.0)
            remv = st.number_input("Removal (%)", min_value=0.0, max_value=100.0, value=54.0, step=1.0)
            remd = st.number_input("Remediation (%)", min_value=0.0, max_value=100.0, value=27.0, step=1.0)

            eprr = {"Engineering": eng, "Preparation": prep, "Removal": remv, "Remediation": remd}
            eprr_total = sum(eprr.values())
            st.caption(f"EPRR total: **{eprr_total:.2f}%**")

            apply_norm = st.checkbox("Normalize EPRR to 100% for this run", value=False)
            if apply_norm and eprr_total > 0 and abs(eprr_total - 100.0) > 1e-6:
                eprr, _ = normalize_to_100(eprr)

        with c2:
            st.markdown("**Financial (%)** (use +/-)")
            sst_pct = st.number_input("SST (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.5)
            owners_pct = st.number_input("Owner's Cost (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.5)
            cont_pct = st.number_input("Contingency (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.5)
            esc_pct = st.number_input("Escalation & Inflation (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.5)

        st.markdown('<h4 style="margin:0;color:#000;">Predict (Single)</h4><p>Step 4</p>', unsafe_allow_html=True)

        project_name = st.text_input("Project Name", placeholder="e.g., Offshore Pipeline Replacement 2025")
        st.caption("Provide feature values (leave blank for NaN).")

        cols_per_row = 3
        new_data = {}
        cols_pred = list(X_processed_pred.columns)
        rows_input = (len(cols_pred) + cols_per_row - 1) // cols_per_row
        for r in range(rows_input):
            row_cols = st.columns(cols_per_row)
            for i in range(cols_per_row):
                idx = r * cols_per_row + i
                if idx < len(cols_pred):
                    col_name = cols_pred[idx]
                    with row_cols[i]:
                        val = st.text_input(col_name, key=f"in_{ds_name_pred}_{col_name}")
                        new_data[col_name] = val

        if st.button("Run Prediction"):
            # Build a single-row DataFrame with the same preprocessing
            input_dict = {}
            for col in cols_pred:
                v = new_data.get(col, "")
                if v is None or str(v).strip() == "":
                    input_dict[col] = np.nan
                else:
                    try:
                        # Try to convert to float; if fails, keep as string
                        input_dict[col] = float(v)
                    except:
                        input_dict[col] = v
            input_df = pd.DataFrame([input_dict])

            # Apply KNN imputation on numeric columns (using the same imputer as on full data)
            if numeric_cols_pred:
                # Refit imputer on full dataset (or we could have stored it)
                full_num = df_pred[numeric_cols_pred]
                imputer_temp = KNNImputer(n_neighbors=5)
                imputer_temp.fit(full_num)
                # Impute the input row for numeric columns
                input_num = imputer_temp.transform(input_df[numeric_cols_pred])
                input_df[numeric_cols_pred] = input_num

            # Get best model name
            best_name = st.session_state.best_model_name_per_dataset.get(ds_name_pred)
            if best_name is None:
                st.error("Please run model training first for this dataset.")
                st.stop()

            # Build pipeline with same preprocessor as during training
            num_cols_train = numeric_cols_pred  # assume same as training (should be stored)
            cat_cols_train = categorical_cols_pred

            numeric_transformer = Pipeline(steps=[
                ('imputer', SimpleImputer(strategy='median')),
                ('scaler', MinMaxScaler())
            ])
            categorical_transformer = Pipeline(steps=[
                ('imputer', SimpleImputer(strategy='constant', fill_value='missing')),
                ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
            ])

            preprocessor = ColumnTransformer(
                transformers=[
                    ('num', numeric_transformer, num_cols_train),
                    ('cat', categorical_transformer, cat_cols_train)
                ]
            )

            ctor = MODEL_CANDIDATES.get(best_name, MODEL_CANDIDATES["RandomForest"])
            try:
                base_model = ctor(42)
            except TypeError:
                base_model = ctor()

            model_pipe = Pipeline(steps=[
                ('preprocessor', preprocessor),
                ('model', base_model)
            ])

            # Train on full dataset
            model_pipe.fit(X_processed_pred, y_processed_pred)

            # Predict
            pred_val = float(model_pipe.predict(input_df)[0])

            owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total = cost_breakdown(
                pred_val, eprr, sst_pct, owners_pct, cont_pct, esc_pct
            )

            result = {"Project Name": project_name}
            result.update({c: new_data.get(c, "") for c in cols_pred})
            result[target_column_pred] = round(pred_val, 2)

            for k, v in eprr_costs.items():
                result[f"{k} Cost"] = v

            result["SST Cost"] = sst_cost
            result["Owner's Cost"] = owners_cost
            result["Cost Contingency"] = contingency_cost
            result["Escalation & Inflation"] = escalation_cost
            result["Grand Total"] = grand_total

            st.session_state.predictions.setdefault(ds_name_pred, []).append(result)
            toast("Prediction added to Results.")

            cA, cB, cC, cD, cE = st.columns(5)
            with cA:
                st.metric("Predicted", f"{currency_pred} {pred_val:,.2f}")
            with cB:
                st.metric("Owner's", f"{currency_pred} {owners_cost:,.2f}")
            with cC:
                st.metric("Contingency", f"{currency_pred} {contingency_cost:,.2f}")
            with cD:
                st.metric("Escalation", f"{currency_pred} {escalation_cost:,.2f}")
            with cE:
                st.metric("Grand Total", f"{currency_pred} {grand_total:,.2f}")

        st.markdown('<h4 style="margin:0;color:#000;">Batch (Excel)</h4>', unsafe_allow_html=True)
        xls = st.file_uploader("Upload Excel for batch prediction", type=["xlsx"], key="batch_xlsx")
        if xls:
            file_id = f"{xls.name}_{xls.size}_{ds_name_pred}"
            if file_id not in st.session_state.processed_excel_files:
                batch_df = pd.read_excel(xls)
                missing = [c for c in X_processed_pred.columns if c not in batch_df.columns]
                if missing:
                    st.error(f"Missing required columns in Excel: {missing}")
                else:
                    # Get best model name
                    best_name = st.session_state.best_model_name_per_dataset.get(ds_name_pred)
                    if best_name is None:
                        st.error("Please run model training first for this dataset.")
                        st.stop()

                    # Prepare batch input: ensure numeric columns are numeric
                    for col in numeric_cols_pred:
                        if col in batch_df.columns:
                            batch_df[col] = pd.to_numeric(batch_df[col], errors='coerce')

                    # Build pipeline
                    num_cols_train = numeric_cols_pred
                    cat_cols_train = categorical_cols_pred

                    numeric_transformer = Pipeline(steps=[
                        ('imputer', SimpleImputer(strategy='median')),
                        ('scaler', MinMaxScaler())
                    ])
                    categorical_transformer = Pipeline(steps=[
                        ('imputer', SimpleImputer(strategy='constant', fill_value='missing')),
                        ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
                    ])

                    preprocessor = ColumnTransformer(
                        transformers=[
                            ('num', numeric_transformer, num_cols_train),
                            ('cat', categorical_transformer, cat_cols_train)
                        ]
                    )

                    ctor = MODEL_CANDIDATES.get(best_name, MODEL_CANDIDATES["RandomForest"])
                    try:
                        base_model = ctor(42)
                    except TypeError:
                        base_model = ctor()

                    model_pipe = Pipeline(steps=[
                        ('preprocessor', preprocessor),
                        ('model', base_model)
                    ])
                    model_pipe.fit(X_processed_pred, y_processed_pred)

                    preds = model_pipe.predict(batch_df[X_processed_pred.columns])

                    for i, row in batch_df.iterrows():
                        name = row.get("Project Name", f"Project {i+1}")

                        owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total = cost_breakdown(
                            float(preds[i]), eprr, sst_pct, owners_pct, cont_pct, esc_pct
                        )

                        entry = {"Project Name": name}
                        entry.update(row[X_processed_pred.columns].to_dict())
                        entry[target_column_pred] = round(float(preds[i]), 2)

                        for k, v in eprr_costs.items():
                            entry[f"{k} Cost"] = v
                        entry["SST Cost"] = sst_cost
                        entry["Owner's Cost"] = owners_cost
                        entry["Cost Contingency"] = contingency_cost
                        entry["Escalation & Inflation"] = escalation_cost
                        entry["Grand Total"] = grand_total

                        st.session_state.predictions.setdefault(ds_name_pred, []).append(entry)

                    st.session_state.processed_excel_files.add(file_id)
                    toast("Batch prediction complete.")
                    st.rerun()
            else:
                st.info("This batch file was already processed (history prevents duplicates).")

    # ========================= RESULTS / EXPORT ==============================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">📄 Results</h3>', unsafe_allow_html=True)

    if not st.session_state.datasets:
        st.info("No dataset. Use the Data section above to upload or load.")
    else:
        ds_name_res = st.selectbox("Dataset", list(st.session_state.datasets.keys()), key="ds_results")
        preds = st.session_state.predictions.get(ds_name_res, [])

        st.markdown(f'<h4 style="margin:0;color:#000;">Project Entries</h4><p>{len(preds)} saved</p>', unsafe_allow_html=True)
        if preds:
            if st.button("🗑️ Delete all entries"):
                st.session_state.predictions[ds_name_res] = []
                to_remove = {fid for fid in st.session_state.processed_excel_files if fid.endswith(ds_name_res)}
                for fid in to_remove:
                    st.session_state.processed_excel_files.remove(fid)
                toast("All entries removed.", "🗑️")
                st.rerun()

        st.markdown('<h4 style="margin:0;color:#000;">Summary Table & Export</h4><p>Download</p>', unsafe_allow_html=True)

        if preds:
            df_preds = pd.DataFrame(preds)
            df_disp = df_preds.copy()
            num_cols = df_disp.select_dtypes(include=[np.number]).columns
            for col in num_cols:
                df_disp[col] = df_disp[col].apply(lambda x: format_with_commas(x))
            st.dataframe(df_disp, use_container_width=True, height=420)

            bio_xlsx = io.BytesIO()
            df_preds.to_excel(bio_xlsx, index=False, engine="openpyxl")
            bio_xlsx.seek(0)

            metrics = st.session_state._last_metrics
            metrics_json = json.dumps(metrics if metrics else {"info": "No metrics"}, indent=2, default=float)

            zip_bio = io.BytesIO()
            with zipfile.ZipFile(zip_bio, "w", zipfile.ZIP_DEFLATED) as zf:
                zf.writestr(f"{ds_name_res}_predictions.xlsx", bio_xlsx.getvalue())
                zf.writestr(f"{ds_name_res}_metrics.json", metrics_json)
            zip_bio.seek(0)

            st.download_button(
                "⬇️ Download All (ZIP)",
                data=zip_bio.getvalue(),
                file_name=f"{ds_name_res}_ABEX_all.zip",
                mime="application/zip",
            )
        else:
            st.info("No data to export yet.")

    # ========================= FLOATER (FPSO/FSO) DECOMMISSIONING ================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">🛳️ Floater (FPSO/FSO) Decommissioning</h3>', unsafe_allow_html=True)

    # -------------------------------------------------------------------------
    # FLOATER TRAINING SECTION
    # -------------------------------------------------------------------------
    st.markdown('<h4 style="margin:0;color:#000;">Floater Training Data</h4><p>Upload CSV with synthetic data</p>', unsafe_allow_html=True)

    floater_file = st.file_uploader(
        "Upload Floater Training CSV", 
        type="csv",
        key="floater_csv_uploader",
        help="CSV should contain columns: UnitType, Location, NoMooringChainAnchor, NoMidWaterArch, MooringHandling, ReimbursableMarkup, NoPipelineRiser, TankCleaning, TankCapacity_bbl, VesselClass, TopsideIsolationCleaning, Number of subsystem, Cost(RM)"
    )

    if floater_file is not None:
        try:
            df_floater = pd.read_csv(floater_file)
            
            # Display basic info
            st.success(f"✅ Loaded {len(df_floater)} rows, {len(df_floater.columns)} columns")
            
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Rows", f"{df_floater.shape[0]:,}")
            with col2:
                st.metric("Columns", f"{df_floater.shape[1]:,}")
            with col3:
                target_col = df_floater.columns[-1] if len(df_floater.columns) > 0 else "None"
                st.metric("Target Column", target_col)
            
            with st.expander("📋 Data Preview (first 10 rows)", expanded=False):
                st.dataframe(df_floater.head(10), use_container_width=True)
            
            with st.expander("📊 Column Details", expanded=False):
                for i, col in enumerate(df_floater.columns):
                    unique_vals = df_floater[col].nunique()
                    dtype = df_floater[col].dtype
                    missing = df_floater[col].isnull().sum()
                    st.write(f"**{i+1}. {col}** - Type: {dtype}, Unique: {unique_vals}, Missing: {missing}")
            
            # Store in session
            st.session_state.floater_dataset = df_floater
            toast("Floater dataset loaded successfully.")
            
        except Exception as e:
            st.error(f"Error loading CSV: {e}")

    # -------------------------------------------------------------------------
    # FLOATER MODEL TRAINING
    # -------------------------------------------------------------------------
    if st.session_state.floater_dataset is not None:
        st.markdown('<h4 style="margin:0;color:#000;">Train Floater Model</h4><p>Step 2: Train prediction model</p>', unsafe_allow_html=True)
        
        df_floater = st.session_state.floater_dataset
        
        # Training parameters
        col1, col2 = st.columns(2)
        with col1:
            test_size = st.slider("Test size", 0.1, 0.3, 0.2, 0.05, key="floater_test_size")
        with col2:
            n_estimators = st.slider("Number of trees", 50, 200, 100, 10, key="floater_n_estimators")
        
        if st.button("🚀 Train Floater Model", key="train_floater_model"):
            with st.spinner("Preprocessing data and training model..."):
                try:
                    # Make a copy for preprocessing
                    df_processed = df_floater.copy()
                    
                    # Handle missing values
                    for col in df_processed.columns:
                        if df_processed[col].dtype == 'object':
                            df_processed[col].fillna(df_processed[col].mode()[0] if not df_processed[col].mode().empty else "Unknown", inplace=True)
                        else:
                            df_processed[col].fillna(df_processed[col].median(), inplace=True)
                    
                    # Prepare features and target
                    X = df_processed.iloc[:, :-1]  # All except last column
                    y = df_processed.iloc[:, -1]   # Last column (Cost)
                    
                    # Store feature columns for later use
                    st.session_state.floater_feature_columns = X.columns.tolist()
                    
                    # Identify categorical and numerical columns
                    categorical_cols = X.select_dtypes(include=['object']).columns.tolist()
                    numerical_cols = X.select_dtypes(include=['int64', 'float64']).columns.tolist()
                    
                    # Create preprocessing pipeline
                    numeric_transformer = Pipeline(steps=[
                        ('imputer', SimpleImputer(strategy='median')),
                        ('scaler', MinMaxScaler())
                    ])
                    
                    categorical_transformer = Pipeline(steps=[
                        ('imputer', SimpleImputer(strategy='constant', fill_value='missing')),
                        ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
                    ])
                    
                    preprocessor = ColumnTransformer(
                        transformers=[
                            ('num', numeric_transformer, numerical_cols),
                            ('cat', categorical_transformer, categorical_cols)
                        ])
                    
                    # Create and train model pipeline
                    model_pipeline = Pipeline([
                        ('preprocessor', preprocessor),
                        ('model', RandomForestRegressor(n_estimators=n_estimators, random_state=42, n_jobs=-1))
                    ])
                    
                    # Split data
                    X_train, X_test, y_train, y_test = train_test_split(
                        X, y, test_size=test_size, random_state=42
                    )
                    
                    # Train model
                    model_pipeline.fit(X_train, y_train)
                    
                    # Evaluate
                    y_pred = model_pipeline.predict(X_test)
                    rmse = np.sqrt(mean_squared_error(y_test, y_pred))
                    r2 = r2_score(y_test, y_pred)
                    
                    # Calculate feature importance if possible
                    feature_importance = None
                    if hasattr(model_pipeline.named_steps['model'], 'feature_importances_'):
                        # Get feature names after preprocessing
                        preprocessor = model_pipeline.named_steps['preprocessor']
                        feature_names = []
                        
                        # Get numerical feature names
                        feature_names.extend(numerical_cols)
                        
                        # Get categorical feature names after one-hot encoding
                        for name, transformer, cols in preprocessor.transformers_:
                            if name == 'cat' and len(cols) > 0:
                                ohe = transformer.named_steps['onehot']
                                feature_names.extend(ohe.get_feature_names_out(cols))
                        
                        importances = model_pipeline.named_steps['model'].feature_importances_
                        
                        # Create feature importance dataframe
                        feature_importance = pd.DataFrame({
                            'Feature': feature_names[:len(importances)],
                            'Importance': importances
                        }).sort_values('Importance', ascending=False).head(15)

                    # Store unique values for categorical columns
                    categorical_unique_values = {}
                    for col in categorical_cols:
                        # Use the processed dataframe (df_processed) which has no missing values
                        unique_vals = df_processed[col].dropna().unique().tolist()
                        categorical_unique_values[col] = unique_vals
                    
                    # Store model and metrics
                    st.session_state.floater_model = {
                        'pipeline': model_pipeline,
                        'metrics': {
                            'rmse': rmse,
                            'r2': r2,
                            'train_samples': len(X_train),
                            'test_samples': len(X_test),
                            'categorical_cols': categorical_cols,
                            'numerical_cols': numerical_cols,
                            'feature_importance': feature_importance
                        },
                        'categorical_unique_values': categorical_unique_values,
                        'feature_columns': X.columns.tolist()
                    }
                    
                    # Show results
                    st.success(f"✅ Model trained successfully!")
                    
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("RMSE", f"RM {rmse:,.2f}")
                    with col2:
                        st.metric("R² Score", f"{r2:.3f}")
                    with col3:
                        st.metric("Training Samples", f"{len(X_train):,}")
                    with col4:
                        st.metric("Test Samples", f"{len(X_test):,}")
                    
                    # Display feature importance if available
                    if feature_importance is not None:
                        st.markdown("##### Top 15 Feature Importances")
                        fig = px.bar(
                            feature_importance,
                            x='Importance',
                            y='Feature',
                            orientation='h',
                            title='Feature Importances',
                            color='Importance',
                            color_continuous_scale='Viridis'
                        )
                        fig.update_layout(height=500)
                        st.plotly_chart(fig, use_container_width=True)
                    
                    # Show model info
                    with st.expander("📊 Model Details", expanded=False):
                        st.write(f"**Model Type:** RandomForestRegressor with {n_estimators} trees")
                        st.write(f"**Categorical Columns:** {', '.join(categorical_cols) if categorical_cols else 'None'}")
                        st.write(f"**Numerical Columns:** {', '.join(numerical_cols) if numerical_cols else 'None'}")
                        
                except Exception as e:
                    st.error(f"Error training model: {e}")
                    import traceback
                    st.error(traceback.format_exc())

    # -------------------------------------------------------------------------
    # FLOATER PREDICTION UI - UPDATED WITH REQUESTED CHANGES
    # -------------------------------------------------------------------------
    if st.session_state.get('floater_model') is not None:
        st.markdown('<h4 style="margin:0;color:#000;">Floater Cost Prediction</h4><p>Step 3: Enter parameters for prediction</p>', unsafe_allow_html=True)
        
        # Get model info
        model_data = st.session_state.floater_model
        model_pipeline = model_data['pipeline']
        metrics = model_data['metrics']
        expected_columns = st.session_state.floater_feature_columns
        categorical_cols = metrics['categorical_cols']
        numerical_cols = metrics['numerical_cols']
        
        # Display model info
        st.info(f"✅ Model ready - R² Score: {metrics['r2']:.3f}, RMSE: RM {metrics['rmse']:,.2f}")
        
        # --- New combined cleaning dropdown (outside form for dynamic visibility) ---
        cleaning_required = st.radio(
            "Cleaning Required?",
            options=["Not required", "Required"],
            index=0,
            key="cleaning_required",
            horizontal=True
        )
        
        # --- Conditionally show Tank Capacity if cleaning is required ---
        tank_capacity = 0.0
        if cleaning_required == "Required":
            tank_capacity = st.number_input(
                "Tank Capacity (bbl)",
                min_value=0,
                max_value=1000000,
                value=400000,
                step=10000,
                key="tank_capacity_input",
                help="Tank capacity in barrels. Required only when cleaning is required."
            )
            # Show vessel class hint
            if tank_capacity > 0:
                if tank_capacity <= 400000:
                    detected_class = "Panamax"
                elif tank_capacity <= 600000:
                    detected_class = "Aframax"
                elif tank_capacity <= 1000000:
                    detected_class = "Suezmax"
                else:
                    detected_class = "VLCC"
                st.caption(f"Detected Vessel Class: {detected_class}")
        
        # --- Form for the remaining inputs ---
        with st.form("floater_prediction_form"):
            st.markdown("**Floater Parameters**")
            
            col1, col2 = st.columns(2)
            
            with col1:
                # UnitType (FPSO/FSO)
                unit_type = st.selectbox(
                    "Unit Type",
                    options=["FPSO", "FSO"],
                    index=0,
                    key="pred_unit_type"
                )
                
                # Location
                location = st.selectbox(
                    "Location",
                    options=["PM", "SB", "SK"],
                    index=0,
                    key="pred_location"
                )
                
                # No Of Mooring Chain And Anchor
                mooring_chain_anchor = st.number_input(
                    "No Of Mooring Chain And Anchor",
                    min_value=0,
                    max_value=20,
                    value=8,
                    step=1,
                    key="pred_mooring_chain"
                )
                
                # No of mid water arch
                mid_water_arch = st.number_input(
                    "No of mid water arch",
                    min_value=0,
                    max_value=10,
                    value=2,
                    step=1,
                    key="pred_mid_water_arch"
                )
                
                # Mooring chain and anchor handling
                mooring_handling = st.selectbox(
                    "Mooring chain and anchor handling",
                    options=[
                        "Cut anchor pile and retrieve chain only",
                        "Mooring Chain and anchor pile leave in situ",
                        "Mooring Chain and drag anchor retrieve/release"
                    ],
                    index=0,
                    key="pred_mooring_handling"
                )
                
                # (Vessel Class removed)
                
            with col2:
                # Reimbursable markup
                markup = st.number_input(
                    "Reimbursable markup (%)",
                    min_value=0.0,
                    max_value=100.0,
                    value=0.0,
                    step=0.5,
                    key="pred_markup"
                )
                
                # No of pipeline/riser
                pipeline_riser = st.number_input(
                    "No of pipeline/riser",
                    min_value=0,
                    max_value=50,
                    value=4,
                    step=1,
                    key="pred_pipeline_riser"
                )
                
                # (Tank cleaning dropdown removed – now handled outside form)
                # (Topside cleaning dropdown removed – same as cleaning)
                
                # (Tank capacity already outside form, will be used later)
                
                # Number of subsystem
                subsystem = st.number_input(
                    "Number of subsystem",
                    min_value=0,
                    max_value=20,
                    value=5,
                    step=1,
                    key="pred_subsystem"
                )
            
            # Prediction button
            submit_prediction = st.form_submit_button("💰 Predict Cost")
            
            if submit_prediction:
                try:
                    # Build input dictionary with correct data types
                    input_dict = {}
                    
                    # Add categorical columns
                    for col in categorical_cols:
                        if col == 'UnitType':
                            input_dict[col] = unit_type
                        elif col == 'Location':
                            input_dict[col] = location
                        elif col == 'MooringHandling':
                            input_dict[col] = mooring_handling
                        elif col == 'TankCleaning':
                            # Use the combined cleaning value
                            input_dict[col] = cleaning_required
                        elif col == 'VesselClass':
                            # Vessel class is inferred from tank capacity; we can compute it here if needed.
                            # The model expects a categorical value; we can map tank capacity to class.
                            # For simplicity, we'll derive it from tank_capacity.
                            if tank_capacity <= 400000:
                                vessel_class = "Panamax"
                            elif tank_capacity <= 600000:
                                vessel_class = "Aframax"
                            elif tank_capacity <= 1000000:
                                vessel_class = "Suezmax"
                            else:
                                vessel_class = "VLCC"
                            input_dict[col] = vessel_class
                        elif col == 'TopsideIsolationCleaning':
                            # Same as tank cleaning
                            input_dict[col] = cleaning_required
                        else:
                            # Default for any other categorical columns
                            input_dict[col] = "Unknown"
                    
                    # Add numerical columns
                    for col in numerical_cols:
                        if col == 'NoMooringChainAnchor':
                            input_dict[col] = float(mooring_chain_anchor)
                        elif col == 'NoMidWaterArch':
                            input_dict[col] = float(mid_water_arch)
                        elif col == 'ReimbursableMarkup':
                            input_dict[col] = float(markup)
                        elif col == 'NoPipelineRiser':
                            input_dict[col] = float(pipeline_riser)
                        elif col == 'TankCapacity_bbl':
                            input_dict[col] = float(tank_capacity)  # may be 0 if cleaning not required
                        elif col == 'Number of subsystem':
                            input_dict[col] = float(subsystem)
                        else:
                            input_dict[col] = 0.0
                    
                    # Create DataFrame with the exact columns the model expects
                    input_df = pd.DataFrame([input_dict])
                    
                    # Ensure column order matches training data
                    input_df = input_df[expected_columns]
                    
                    # Make prediction
                    prediction = model_pipeline.predict(input_df)[0]
                    
                    # Store prediction with user-friendly column names
                    prediction_record = {
                        'Unit Type': unit_type,
                        'Location': location,
                        'Mooring Chain & Anchor': mooring_chain_anchor,
                        'Mid Water Arch': mid_water_arch,
                        'Mooring Handling': mooring_handling,
                        'Markup (%)': markup,
                        'Pipeline/Riser': pipeline_riser,
                        'Tank Cleaning': cleaning_required,
                        'Tank Capacity (bbl)': tank_capacity,
                        'Vessel Class': input_dict.get('VesselClass', 'N/A'),
                        'Topside Cleaning': cleaning_required,
                        'Subsystems': subsystem,
                        'Predicted Cost (RM)': float(prediction),
                        'Timestamp': pd.Timestamp.now().strftime("%Y-%m-%d %H:%M:%S")
                    }
                    
                    st.session_state.floater_predictions.append(prediction_record)
                    
                    # Display results
                    st.success(f"✅ Cost Prediction Complete!")
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric(
                            "Base Cost",
                            f"RM {prediction:,.2f}"
                        )
                    
                    with col2:
                        # Calculate with markup
                        markup_amount = prediction * (markup / 100)
                        st.metric(
                            "Markup Amount",
                            f"RM {markup_amount:,.2f}",
                            delta=f"{markup}%"
                        )
                    
                    with col3:
                        total_with_markup = prediction + markup_amount
                        st.metric(
                            "Total Cost (with markup)",
                            f"RM {total_with_markup:,.2f}"
                        )
                    
                    # Show input summary
                    with st.expander("📋 Input Summary", expanded=False):
                        summary_df = pd.DataFrame([prediction_record])
                        # Drop timestamp for cleaner display
                        display_cols = [c for c in summary_df.columns if c != 'Timestamp']
                        st.dataframe(summary_df[display_cols], use_container_width=True)
                    
                    # Show confidence interval
                    st.info(f"📊 Model confidence: R² = {metrics['r2']:.3f}, Typical error range: ±RM {metrics['rmse']:,.2f}")
                    
                except Exception as e:
                    st.error(f"Error making prediction: {e}")
                    import traceback
                    st.error(traceback.format_exc())

    # -------------------------------------------------------------------------
    # FLOATER PREDICTION HISTORY
    # -------------------------------------------------------------------------
    if st.session_state.floater_predictions:
        st.markdown('<h4 style="margin:0;color:#000;">Prediction History</h4><p>Previous predictions</p>', unsafe_allow_html=True)
        
        # Convert to DataFrame for display
        history_df = pd.DataFrame(st.session_state.floater_predictions)
        
        # Display in a nice table
        display_cols = ['Timestamp', 'Unit Type', 'Location', 'Tank Cleaning', 'Predicted Cost (RM)']
        st.dataframe(
            history_df[display_cols].sort_values('Timestamp', ascending=False),
            use_container_width=True
        )
        
        # Export options
        col1, col2 = st.columns(2)
        
        with col1:
            # Export to Excel
            if not history_df.empty:
                excel_buffer = io.BytesIO()
                with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                    history_df.to_excel(writer, index=False, sheet_name='Predictions')
                excel_buffer.seek(0)
                
                st.download_button(
                    label="📥 Download Predictions (Excel)",
                    data=excel_buffer,
                    file_name="floater_predictions.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
        
        with col2:
            # Clear history button
            if st.button("🧹 Clear Prediction History", key="clear_floater_history"):
                st.session_state.floater_predictions = []
                toast("Prediction history cleared.")
                st.rerun()

    # -------------------------------------------------------------------------
    # FLOATER BATCH PREDICTION - FIXED TO MATCH SINGLE PREDICTION
    # -------------------------------------------------------------------------
    if st.session_state.get('floater_model') is not None:
        st.markdown('<h4 style="margin:0;color:#000;">Batch Prediction</h4><p>Upload Excel for multiple predictions</p>', unsafe_allow_html=True)

        # --- Template download button ---
        if st.button("📥 Download Batch Template", key="floater_template_btn"):
            # Note: Template now reflects combined cleaning and no vessel class column.
            # The user should still include both TankCleaning and TopsideIsolationCleaning columns (same value).
            # VesselClass can be omitted; it will be derived from TankCapacity during batch processing.
            template_data = {
                'UnitType': ['FPSO'],
                'Location': ['PM'],
                'NoMooringChainAnchor': [8],
                'NoMidWaterArch': [2],
                'MooringHandling': ['Cut anchor pile and retrieve chain only'],
                'ReimbursableMarkup': [0.0],
                'NoPipelineRiser': [4],
                'TankCleaning': ['Not required'],
                'TankCapacity_bbl': [0],  # 0 when not required
                'TopsideIsolationCleaning': ['Not required'],
                'Number of subsystem': [5],
            }
            template_df = pd.DataFrame(template_data)
            
            output = io.BytesIO()
            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                template_df.to_excel(writer, index=False, sheet_name='Template')
                workbook = writer.book
                worksheet = writer.sheets['Template']
                cat_options = {
                    'UnitType': ['FPSO', 'FSO'],
                    'Location': ['PM', 'SB', 'SK'],
                    'MooringHandling': [
                        'Cut anchor pile and retrieve chain only',
                        'Mooring Chain and anchor pile leave in situ',
                        'Mooring Chain and drag anchor retrieve/release'
                    ],
                    'TankCleaning': ['Required', 'Not required'],
                    'TopsideIsolationCleaning': ['Required', 'Not required']
                }
                col_mapping = {name: idx+1 for idx, name in enumerate(template_df.columns)}
                for col_name, options in cat_options.items():
                    if col_name in col_mapping:
                        col_letter = get_column_letter(col_mapping[col_name])
                        dv = DataValidation(type="list", formula1=f'"{",".join(options)}"')
                        worksheet.add_data_validation(dv)
                        dv.add(f'{col_letter}2:{col_letter}1048576')
            output.seek(0)
            st.download_button(
                "⬇️ Download Excel Template",
                data=output,
                file_name="floater_prediction_template.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="floater_template_download"
            )

        # Get model info
        model_data = st.session_state.floater_model
        model_pipeline = model_data['pipeline']
        metrics = model_data['metrics']
        expected_columns = st.session_state.floater_feature_columns
        categorical_cols = metrics['categorical_cols']
        numerical_cols = metrics['numerical_cols']
        cat_unique = model_data.get('categorical_unique_values', {})

        batch_file = st.file_uploader(
            "Upload Excel file for batch prediction",
            type=["xlsx", "xls"],
            key="floater_batch_file"
        )

        if batch_file is not None:
            try:
                batch_df = pd.read_excel(batch_file)
                st.write("📊 Uploaded Data Preview:")
                st.dataframe(batch_df.head())

                missing_cols = [col for col in expected_columns if col not in batch_df.columns]
                if missing_cols:
                    st.error(f"Missing required columns: {', '.join(missing_cols)}")
                    st.info(f"Expected columns: {', '.join(expected_columns)}")
                else:
                    if st.button("🔢 Run Batch Prediction", key="run_batch_floater"):
                        with st.spinner("Processing batch prediction..."):
                            try:
                                processed_rows = []
                                for idx, row in batch_df.iterrows():
                                    input_dict = {}

                                    # ---- Categorical columns ----
                                    for col in categorical_cols:
                                        raw_val = row[col]
                                        if pd.isna(raw_val) or raw_val is None:
                                            input_dict[col] = "Unknown"
                                        else:
                                            val = str(raw_val).strip()
                                            known_values = cat_unique.get(col, [])
                                            
                                            # Define mapping function for boolean-like values
                                            def map_to_known(v, known):
                                                if not known:
                                                    return v
                                                if v in known:
                                                    return v
                                                # case-insensitive
                                                for k in known:
                                                    if k.lower() == v.lower():
                                                        return k
                                                # map common boolean representations
                                                v_low = v.lower()
                                                if v_low in ['1', 'true', 'yes']:
                                                    # pick the one that likely means "required" (not containing "not")
                                                    for k in known:
                                                        if 'not' not in k.lower():
                                                            return k
                                                if v_low in ['0', 'false', 'no']:
                                                    for k in known:
                                                        if 'not' in k.lower():
                                                            return k
                                                # fallback
                                                return known[0]
                                            
                                            mapped_val = map_to_known(val, known_values)
                                            if mapped_val != val:
                                                st.warning(f"Row {idx+1}: Column '{col}' value '{raw_val}' mapped to '{mapped_val}'.")
                                            input_dict[col] = mapped_val

                                    # ---- Numerical columns ----
                                    for col in numerical_cols:
                                        raw_val = row[col]
                                        if pd.isna(raw_val) or raw_val is None:
                                            input_dict[col] = 0.0
                                        else:
                                            # Remove commas, percent signs, and convert
                                            try:
                                                # Handle strings like "1,000" or "10%"
                                                cleaned = str(raw_val).replace(',', '').replace('%', '').strip()
                                                input_dict[col] = float(cleaned)
                                            except:
                                                st.warning(f"Row {idx+1}: Column '{col}' value '{raw_val}' could not be parsed as number. Using 0.")
                                                input_dict[col] = 0.0

                                    processed_rows.append(input_dict)

                                # Create DataFrame with correct column order
                                processed_df = pd.DataFrame(processed_rows)
                                processed_df = processed_df[expected_columns]

                                # Force numerical columns to float (safety)
                                for col in numerical_cols:
                                    if col in processed_df.columns:
                                        processed_df[col] = pd.to_numeric(processed_df[col], errors='coerce').fillna(0.0)

                                # Make predictions
                                predictions = model_pipeline.predict(processed_df)

                                # Add predictions to original dataframe
                                batch_df['Predicted Cost (RM)'] = predictions
                                if 'ReimbursableMarkup' in batch_df.columns:
                                    batch_df['Total with Markup (RM)'] = batch_df.apply(
                                        lambda row: row['Predicted Cost (RM)'] * (1 + row['ReimbursableMarkup']/100),
                                        axis=1
                                    )

                                st.success(f"✅ Batch prediction complete! Processed {len(batch_df)} records.")
                                col1, col2, col3 = st.columns(3)
                                with col1:
                                    st.metric("Min Prediction", f"RM {predictions.min():,.2f}")
                                with col2:
                                    st.metric("Max Prediction", f"RM {predictions.max():,.2f}")
                                with col3:
                                    st.metric("Avg Prediction", f"RM {predictions.mean():,.2f}")

                                with st.expander("📊 Detailed Results", expanded=True):
                                    st.dataframe(batch_df, use_container_width=True)

                                # Download button
                                batch_buffer = io.BytesIO()
                                with pd.ExcelWriter(batch_buffer, engine='openpyxl') as writer:
                                    batch_df.to_excel(writer, index=False, sheet_name='Batch Predictions')
                                batch_buffer.seek(0)
                                st.download_button(
                                    label="📥 Download Batch Results",
                                    data=batch_buffer,
                                    file_name="floater_batch_predictions.xlsx",
                                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                                )

                                # Compare with latest single prediction
                                if st.session_state.floater_predictions:
                                    latest_single = st.session_state.floater_predictions[-1]
                                    st.info(f"Latest single prediction: RM {latest_single['Predicted Cost (RM)']:,.2f}")

                            except Exception as e:
                                st.error(f"Error during prediction: {e}")
                                import traceback
                                st.error(traceback.format_exc())

            except Exception as e:
                st.error(f"Error processing batch file: {e}")

    # -------------------------------------------------------------------------
    # RESET FLOATER SECTION
    # -------------------------------------------------------------------------
    st.markdown("---")
    if st.button("🔄 Reset Floater Section", key="reset_floater"):
        st.session_state.floater_dataset = None
        st.session_state.floater_model = None
        st.session_state.floater_predictions = []
        st.session_state.floater_feature_columns = []
        toast("Floater section reset successfully.")
        st.rerun()


# =======================================================================================
# PROJECT BUILDER TAB
# =======================================================================================
with tab_pb:
    st.markdown('<h4 style="margin:0;color:#000;">Project Builder</h4><p>Assemble multi-component ABEX projects</p>', unsafe_allow_html=True)

    if not st.session_state.datasets:
        st.info("No dataset. Go to **Data** tab to upload or load.")
    else:
        colA, colB = st.columns([2, 1])
        with colA:
            new_project_name = st.text_input("New Project Name", placeholder="e.g., ABEX 2026", key="pb_new_project_name")
        with colB:
            if new_project_name and new_project_name not in st.session_state.projects:
                if st.button("Create Project", key="pb_create_project_btn"):
                    st.session_state.projects[new_project_name] = {"components": [], "totals": {}, "currency": ""}
                    toast(f"Project '{new_project_name}' created.")
                    st.rerun()

        if not st.session_state.projects:
            st.info("Create a project above, then add components.")
        else:
            existing_projects = list(st.session_state.projects.keys())
            proj_sel = st.selectbox("Select project to work on", existing_projects, key="pb_project_select")

            ds_names = sorted(st.session_state.datasets.keys())
            dataset_for_comp = st.selectbox("Dataset for this component", ds_names, key="pb_dataset_for_component")

            df_comp = st.session_state.datasets[dataset_for_comp]
            currency_ds = get_currency_symbol(df_comp)

            # Prepare data: impute numeric only, keep categorical
            num_cols_comp = df_comp.select_dtypes(include=np.number).columns.tolist()
            cat_cols_comp = df_comp.select_dtypes(include='object').columns.tolist()
            
            if num_cols_comp:
                imputer_num = KNNImputer(n_neighbors=5)
                imputed_num = imputer_num.fit_transform(df_comp[num_cols_comp])
                df_imputed_num = pd.DataFrame(imputed_num, columns=num_cols_comp, index=df_comp.index)
            else:
                df_imputed_num = pd.DataFrame()
            
            df_cat_comp = df_comp[cat_cols_comp] if cat_cols_comp else pd.DataFrame()
            df_imputed_comp = pd.concat([df_imputed_num, df_cat_comp], axis=1)
            df_imputed_comp = df_imputed_comp[df_comp.columns]
            
            X_comp = df_imputed_comp.iloc[:, :-1]
            y_comp = df_imputed_comp.iloc[:, -1]
            target_column_comp = y_comp.name

            default_label = st.session_state.component_labels.get(dataset_for_comp, "")
            component_type = st.text_input(
                "Component type (Asset / Scope)",
                value=(default_label or "Platform / Pipeline / Subsea / Well"),
                key=f"pb_component_type_{proj_sel}",
            )

            st.markdown("**Component Feature Inputs**")
            feat_cols = list(X_comp.columns)
            comp_inputs = {}
            cols_per_row = 2
            rows_input = (len(feat_cols) + cols_per_row - 1) // cols_per_row
            for r in range(rows_input):
                row_cols = st.columns(cols_per_row)
                for i in range(cols_per_row):
                    idx = r * cols_per_row + i
                    if idx < len(feat_cols):
                        col_name = feat_cols[idx]
                        with row_cols[i]:
                            key = f"pb_{proj_sel}_{dataset_for_comp}_feat_{col_name}"
                            comp_inputs[col_name] = st.text_input(col_name, key=key)

            st.markdown("---")
            st.markdown("**Cost Percentage Inputs**")
            cp1, cp2 = st.columns(2)
            with cp1:
                st.markdown("EPRR (%) — use +/-")
                eng_pb = st.number_input("Engineering", 0.0, 100.0, 12.0, 1.0, key=f"pb_eng_{proj_sel}")
                prep_pb = st.number_input("Preparation", 0.0, 100.0, 7.0, 1.0, key=f"pb_prep_{proj_sel}")
                remv_pb = st.number_input("Removal", 0.0, 100.0, 54.0, 1.0, key=f"pb_remv_{proj_sel}")
                remd_pb = st.number_input("Remediation", 0.0, 100.0, 27.0, 1.0, key=f"pb_remd_{proj_sel}")

                eprr_pb = {"Engineering": eng_pb, "Preparation": prep_pb, "Removal": remv_pb, "Remediation": remd_pb}
                eprr_total_pb = sum(eprr_pb.values())
                st.caption(f"EPRR total: **{eprr_total_pb:.2f}%**")
                apply_norm_pb = st.checkbox("Normalize to 100% for this component", value=False, key=f"pb_norm_{proj_sel}")
                if apply_norm_pb and eprr_total_pb > 0 and abs(eprr_total_pb - 100.0) > 1e-6:
                    eprr_pb, _ = normalize_to_100(eprr_pb)

            with cp2:
                st.markdown("Financial (%) — use +/-")
                sst_pb = st.number_input("SST", 0.0, 100.0, 0.0, 0.5, key=f"pb_sst_{proj_sel}")
                owners_pb = st.number_input("Owner's Cost", 0.0, 100.0, 0.0, 0.5, key=f"pb_owners_{proj_sel}")
                cont_pb = st.number_input("Contingency", 0.0, 100.0, 0.0, 0.5, key=f"pb_cont_{proj_sel}")
                esc_pb = st.number_input("Escalation & Inflation", 0.0, 100.0, 0.0, 0.5, key=f"pb_esc_{proj_sel}")

            if st.button("➕ Predict & Add Component", key=f"pb_add_comp_{proj_sel}_{dataset_for_comp}"):
                row_payload = {}
                for f in feat_cols:
                    v = comp_inputs.get(f, "")
                    if v is None or str(v).strip() == "":
                        row_payload[f] = np.nan
                    else:
                        try:
                            row_payload[f] = float(v)
                        except Exception:
                            row_payload[f] = np.nan

                try:
                    # Get best model name
                    best_name = st.session_state.best_model_name_per_dataset.get(dataset_for_comp)
                    if best_name is None:
                        st.error("Please run model training first for this dataset.")
                        st.stop()

                    # Build pipeline with same preprocessor as during training
                    num_cols_train = X_comp.select_dtypes(include=np.number).columns.tolist()
                    cat_cols_train = X_comp.select_dtypes(include='object').columns.tolist()

                    numeric_transformer = Pipeline(steps=[
                        ('imputer', SimpleImputer(strategy='median')),
                        ('scaler', MinMaxScaler())
                    ])
                    categorical_transformer = Pipeline(steps=[
                        ('imputer', SimpleImputer(strategy='constant', fill_value='missing')),
                        ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
                    ])

                    preprocessor = ColumnTransformer(
                        transformers=[
                            ('num', numeric_transformer, num_cols_train),
                            ('cat', categorical_transformer, cat_cols_train)
                        ]
                    )

                    ctor = MODEL_CANDIDATES.get(best_name, MODEL_CANDIDATES["RandomForest"])
                    try:
                        base_model = ctor(42)
                    except TypeError:
                        base_model = ctor()

                    model_pipe = Pipeline(steps=[
                        ('preprocessor', preprocessor),
                        ('model', base_model)
                    ])

                    # Train on full dataset
                    model_pipe.fit(X_comp, y_comp)

                    # Prepare input row
                    input_df = pd.DataFrame([row_payload])
                    # Ensure numeric columns are float
                    for col in num_cols_train:
                        if col in input_df.columns:
                            input_df[col] = pd.to_numeric(input_df[col], errors='coerce')

                    base_pred = float(model_pipe.predict(input_df)[0])

                    owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total = cost_breakdown(
                        base_pred, eprr_pb, sst_pb, owners_pb, cont_pb, esc_pb
                    )

                    comp_entry = {
                        "component_type": component_type or default_label or "Component",
                        "dataset": dataset_for_comp,
                        "model_used": best_name,
                        "inputs": {k: row_payload[k] for k in feat_cols},
                        "prediction": base_pred,
                        "breakdown": {
                            "eprr_costs": eprr_costs,
                            "eprr_pct": eprr_pb,
                            "sst_cost": sst_cost,
                            "owners_cost": owners_cost,
                            "contingency_cost": contingency_cost,
                            "escalation_cost": escalation_cost,
                            "grand_total": grand_total,
                            "target_col": target_column_comp,
                        },
                    }

                    st.session_state.projects[proj_sel]["components"].append(comp_entry)
                    st.session_state.component_labels[dataset_for_comp] = component_type or default_label

                    # Currency for project should follow dataset currency
                    st.session_state.projects[proj_sel]["currency"] = currency_ds

                    toast(f"Component added to project '{proj_sel}'.")
                    st.rerun()
                except Exception as e:
                    st.error(f"Failed to predict component ABEX: {e}")

            st.markdown("---")
            st.markdown("### Current Project Overview")

            proj = st.session_state.projects[proj_sel]
            comps = proj.get("components", [])
            if not comps:
                st.info("No components yet. Add at least one above.")
            else:
                dfc = project_components_df(proj)
                curr = proj.get("currency", "") or currency_ds

                st.dataframe(
                    dfc.style.format({"Base ABEX": "{:,.2f}", "Grand Total": "{:,.2f}"}),
                    use_container_width=True,
                )

                total_ABEX = float(dfc["Base ABEX"].sum())
                total_grand = float(dfc["Grand Total"].sum())
                proj["totals"] = {"ABEX_sum": total_ABEX, "grand_total": total_grand}

                col_t1, col_t2 = st.columns(2)
                with col_t1:
                    st.metric("Project ABEX", f"{curr} {total_ABEX:,.2f}")
                with col_t2:
                    st.metric("Project Grand Total", f"{curr} {total_grand:,.2f}")

                st.markdown("#### Component Cost Composition")
                comp_cost_rows = []
                for c in comps:
                    comp_cost_rows.append(
                        {
                            "Component": c["component_type"],
                            "ABEX": float(c["prediction"]),
                            "Owner": float(c["breakdown"]["owners_cost"]),
                            "Contingency": float(c["breakdown"]["contingency_cost"]),
                            "Escalation": float(c["breakdown"]["escalation_cost"]),
                            "SST": float(c["breakdown"]["sst_cost"]),
                        }
                    )
                df_cost = pd.DataFrame(comp_cost_rows)
                if not df_cost.empty:
                    df_melt = df_cost.melt(id_vars="Component", var_name="Cost Type", value_name="Value")
                    fig_stack = px.bar(
                        df_melt,
                        x="Component",
                        y="Value",
                        color="Cost Type",
                        barmode="stack",
                        labels={"Value": f"Cost ({curr})"},
                    )
                    st.plotly_chart(fig_stack, use_container_width=True)

                st.markdown("#### Components")
                for idx, c in enumerate(comps):
                    col1, col2, col3 = st.columns([4, 2, 1])
                    with col1:
                        st.write(f"**{c['component_type']}** — *{c['dataset']}* — {c.get('model_used', 'N/A')}")
                    with col2:
                        st.write(f"Grand Total: {curr} {c['breakdown']['grand_total']:,.2f}")
                    with col3:
                        if st.button("🗑️", key=f"pb_del_comp_{proj_sel}_{idx}"):
                            comps.pop(idx)
                            toast("Component removed.", "🗑️")
                            st.rerun()

                st.markdown("---")
                st.markdown("#### Export / Import Project")

                col_dl1, col_dl2, col_dl3 = st.columns(3)

                with col_dl1:
                    excel_report = create_project_excel_report_ABEX(proj_sel, proj, curr)
                    st.download_button(
                        "⬇️ Download Project Excel",
                        data=excel_report,
                        file_name=f"{proj_sel}_ABEX_Report.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    )

                with col_dl2:
                    pptx_report = create_project_pptx_report_ABEX(proj_sel, proj, curr)
                    st.download_button(
                        "⬇️ Download Project PowerPoint",
                        data=pptx_report,
                        file_name=f"{proj_sel}_ABEX_Report.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )

                with col_dl3:
                    st.download_button(
                        "⬇️ Download Project (JSON)",
                        data=json.dumps(proj, indent=2, default=float),
                        file_name=f"{proj_sel}.json",
                        mime="application/json",
                    )

                up_json = st.file_uploader("Import project JSON", type=["json"], key=f"pb_import_{proj_sel}")
                if up_json is not None:
                    try:
                        data = json.load(up_json)
                        st.session_state.projects[proj_sel] = data
                        toast("Project imported successfully.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Failed to import project JSON: {e}")


# =======================================================================================
# COMPARE PROJECTS TAB
# =======================================================================================
with tab_compare:
    st.markdown('<h4 style="margin:0;color:#000;">Compare Projects</h4><p>Portfolio-level ABEX view</p>', unsafe_allow_html=True)

    proj_names = list(st.session_state.projects.keys())
    if len(proj_names) < 2:
        st.info("Create at least two projects in the Project Builder tab to compare.")
    else:
        compare_sel = st.multiselect(
            "Select projects to compare", proj_names, default=proj_names[:2], key="compare_projects_sel"
        )

        if len(compare_sel) < 2:
            st.warning("Select at least two projects for a meaningful comparison.")
        else:
            rows = []
            for p in compare_sel:
                proj = st.session_state.projects[p]
                dfc = project_components_df(proj)

                ABEX = float(dfc["Base ABEX"].sum()) if not dfc.empty else 0.0
                owners = float(dfc["Owner's Cost"].sum()) if not dfc.empty else 0.0
                cont = float(dfc["Contingency"].sum()) if not dfc.empty else 0.0
                esc = float(dfc["Escalation"].sum()) if not dfc.empty else 0.0
                sst = float(dfc["SST"].sum()) if not dfc.empty else 0.0
                grand_total = float(dfc["Grand Total"].sum()) if not dfc.empty else 0.0

                proj["totals"] = {"ABEX_sum": ABEX, "grand_total": grand_total}
                rows.append(
                    {
                        "Project": p,
                        "Components": len(proj.get("components", [])),
                        "ABEX Sum": ABEX,
                        "Owner": owners,
                        "Contingency": cont,
                        "Escalation": esc,
                        "SST": sst,
                        "Grand Total": grand_total,
                        "Currency": proj.get("currency", ""),
                    }
                )

            df_proj = pd.DataFrame(rows)

            st.dataframe(
                df_proj[["Project", "Components", "ABEX Sum", "Grand Total"]].style.format(
                    {"ABEX Sum": "{:,.2f}", "Grand Total": "{:,.2f}"}
                ),
                use_container_width=True,
            )

            st.markdown("#### Grand Total by Project")
            fig_gt = px.bar(df_proj, x="Project", y="Grand Total", text="Grand Total", barmode="group")
            fig_gt.update_traces(texttemplate="%{text:,.0f}", textposition="outside")
            st.plotly_chart(fig_gt, use_container_width=True)

            st.markdown("#### Stacked Cost Composition by Project")
            df_melt = df_proj.melt(
                id_vars=["Project"],
                value_vars=["ABEX Sum", "Owner", "Contingency", "Escalation", "SST"],
                var_name="Cost Type",
                value_name="Value",
            )
            fig_comp = px.bar(df_melt, x="Project", y="Value", color="Cost Type", barmode="stack")
            st.plotly_chart(fig_comp, use_container_width=True)

            st.markdown("#### Component-Level Details")
            for p in compare_sel:
                proj = st.session_state.projects[p]
                comps = proj.get("components", [])
                if not comps:
                    continue
                with st.expander(f"Project: {p}"):
                    rows_c = []
                    for c in comps:
                        eprr_costs = c["breakdown"].get("eprr_costs", {})
                        eprr_str = ", ".join(f"{k}: {v:,.0f}" for k, v in eprr_costs.items() if float(v) != 0)
                        rows_c.append(
                            {
                                "Component": c["component_type"],
                                "Dataset": c["dataset"],
                                "Base ABEX": c["prediction"],
                                "Owner": c["breakdown"]["owners_cost"],
                                "Contingency": c["breakdown"]["contingency_cost"],
                                "Escalation": c["breakdown"]["escalation_cost"],
                                "SST": c["breakdown"]["sst_cost"],
                                "Grand Total": c["breakdown"]["grand_total"],
                                "EPRR Costs": eprr_str,
                            }
                        )
                    df_compd = pd.DataFrame(rows_c)
                    st.dataframe(
                        df_compd.style.format(
                            {
                                "Base ABEX": "{:,.2f}",
                                "Owner": "{:,.2f}",
                                "Contingency": "{:,.2f}",
                                "Escalation": "{:,.2f}",
                                "SST": "{:,.2f}",
                                "Grand Total": "{:,.2f}",
                            }
                        ),
                        use_container_width=True,
                    )

            st.markdown("---")
            st.markdown("#### Download Comparison Reports")

            col_c1, col_c2 = st.columns(2)
            projects_to_export = {name: st.session_state.projects[name] for name in compare_sel}
            currency_comp = st.session_state.projects[compare_sel[0]].get("currency", "")

            with col_c1:
                excel_comp = create_comparison_excel_report_ABEX(projects_to_export, currency_comp)
                st.download_button(
                    "⬇️ Download Comparison Excel",
                    data=excel_comp,
                    file_name="ABEX_Projects_Comparison.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )

            with col_c2:
                pptx_comp = create_comparison_pptx_report_ABEX(projects_to_export, currency_comp)
                st.download_button(
                    "⬇️ Download Comparison PowerPoint",
                    data=pptx_comp,
                    file_name="ABEX_Projects_Comparison.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
