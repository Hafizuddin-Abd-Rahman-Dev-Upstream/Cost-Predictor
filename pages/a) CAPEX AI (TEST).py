# ======================================================================================
# CAPEX AI RT2026 —  STREAMLIT TEST VERSION
# - Data (upload/load GitHub)
# - Train + Predict (single + batch)
# - Project Builder
# - 🎲 Monte Carlo (NEW top-level tab)
# - Compare Projects
# - Exports (ZIP, Excel, PPTX)


import io
import json
import zipfile
import re
import requests
import numpy as np
import pandas as pd
import streamlit as st

# ---------------------------
# HARD GUARD: sklearn missing
# ---------------------------
try:
    from sklearn.impute import KNNImputer, SimpleImputer
    from sklearn.model_selection import train_test_split
    from sklearn.preprocessing import MinMaxScaler
    from sklearn.ensemble import RandomForestRegressor, GradientBoostingRegressor
    from sklearn.linear_model import Ridge, Lasso
    from sklearn.svm import SVR
    from sklearn.tree import DecisionTreeRegressor
    from sklearn.pipeline import Pipeline
    from sklearn.metrics import mean_squared_error, r2_score
except Exception as e:
    st.set_page_config(page_title="CAPEX AI RT2026", page_icon="💠", layout="wide")
    st.error("❌ Missing dependency: scikit-learn (sklearn).")
    st.code(
        "Create requirements.txt in your repo root with:\n\n"
        "streamlit\npandas\nnumpy\nscipy\nscikit-learn\nplotly\nmatplotlib\npython-pptx\nopenpyxl\nrequests\n"
    )
    st.info("Then redeploy/reboot the Streamlit Cloud app (Manage app → Reboot).")
    st.exception(e)
    st.stop()

from scipy.stats import linregress

import plotly.express as px
import plotly.graph_objects as go

import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.utils import get_column_letter


# ---------------------------------------------------------------------------------------
# PAGE CONFIG
# ---------------------------------------------------------------------------------------
st.set_page_config(
    page_title="CAPEX AI RT2026",
    page_icon="💠",
    layout="wide",
    initial_sidebar_state="expanded",
)


# ---------------------------------------------------------------------------------------
# SESSION STATE
# ---------------------------------------------------------------------------------------
def _init_state():
    defaults = {
        "authenticated": False,
        "datasets": {},
        "predictions": {},
        "processed_excel_files": set(),
        "_last_metrics": None,
        "projects": {},
        "component_labels": {},
        "uploader_nonce": 0,
        "_widget_nonce": 0,  # used to force uniqueness when needed
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v


_init_state()


# ---------------------------------------------------------------------------------------
# THEME / CSS (optional)
# ---------------------------------------------------------------------------------------
PETRONAS = {
    "teal": "#00A19B",
    "teal_dark": "#008C87",
    "purple": "#6C4DD3",
    "white": "#FFFFFF",
    "black": "#0E1116",
    "border": "rgba(0,0,0,0.10)",
}

st.markdown(
    f"""
<style>
html, body {{ font-family: sans-serif; }}
[data-testid="stAppViewContainer"] {{
  background: {PETRONAS["white"]};
  color: {PETRONAS["black"]};
}}
#MainMenu, footer {{ visibility: hidden; }}

.petronas-hero {{
  border-radius: 16px;
  padding: 18px 22px;
  margin: 6px 0 18px 0;
  color: #fff;
  background: linear-gradient(135deg, {PETRONAS["teal"]}, {PETRONAS["purple"]}, {PETRONAS["black"]});
  box-shadow: 0 10px 24px rgba(0,0,0,.12);
}}
.petronas-hero h1 {{ margin: 0 0 4px; font-weight: 800; }}
.petronas-hero p {{ margin: 0; opacity: .9; font-weight: 500; }}
</style>
""",
    unsafe_allow_html=True,
)

st.markdown(
    """
<div class="petronas-hero">
  <h1>CAPEX AI RT2026 (TEST)</h1>
  <p>Data-driven CAPEX prediction + Project Builder + Monte Carlo</p>
</div>
""",
    unsafe_allow_html=True,
)


# ---------------------------------------------------------------------------------------
# AUTH (optional)
# ---------------------------------------------------------------------------------------
# If you DON'T want login on test, set ENABLE_AUTH = False
ENABLE_AUTH = True

if ENABLE_AUTH:
    APPROVED_EMAILS = [str(e).strip().lower() for e in st.secrets.get("emails", [])]
    correct_password = st.secrets.get("password", None)

    if not st.session_state.authenticated:
        with st.form("login_form"):
            st.markdown("#### 🔐 Access Required")
            email = (st.text_input("Email Address", key="auth_email") or "").strip().lower()
            password = st.text_input("Access Password", type="password", key="auth_pw")
            submitted = st.form_submit_button("Login")

            if submitted:
                ok = (email in APPROVED_EMAILS) and (password == correct_password)
                if ok:
                    st.session_state.authenticated = True
                    st.success("✅ Access granted.")
                    st.rerun()
                else:
                    st.error("❌ Invalid credentials.")
        st.stop()


# ---------------------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------------------
def toast(msg, icon="✅"):
    try:
        st.toast(f"{icon} {msg}")
    except Exception:
        st.success(f"{icon} {msg}")


def format_with_commas(num):
    try:
        return f"{float(num):,.2f}"
    except Exception:
        return str(num)


def normalize_to_100(d: dict):
    total = sum(float(v) for v in d.values())
    if total <= 0:
        return d, total
    out = {k: float(v) * 100.0 / total for k, v in d.items()}
    keys = list(out.keys())
    rounded = {k: round(out[k], 2) for k in keys}
    diff = 100.0 - sum(rounded.values())
    if keys:
        rounded[keys[-1]] = round(rounded[keys[-1]] + diff, 2)
    return rounded, total


def is_junk_col(colname: str) -> bool:
    h = str(colname).strip().upper()
    return (not h) or h.startswith("UNNAMED") or h in {"INDEX", "IDX"}


def currency_from_header(header: str) -> str:
    h = (header or "").strip().upper()
    if "€" in h:
        return "€"
    if "£" in h:
        return "£"
    if "$" in h:
        return "$"
    if re.search(r"\bUSD\b", h):
        return "USD"
    if re.search(r"\b(MYR|RM)\b", h):
        return "RM"
    return ""


def get_currency_symbol(df: pd.DataFrame, target_col: str | None = None) -> str:
    if df is None or df.empty:
        return ""
    if target_col and target_col in df.columns:
        return currency_from_header(str(target_col))
    for c in reversed(df.columns):
        if not is_junk_col(c):
            sym = currency_from_header(str(c))
            if sym:
                return sym
    return ""


def cost_breakdown(
    base_pred: float,
    eprr: dict,
    sst_pct: float,
    owners_pct: float,
    cont_pct: float,
    esc_pct: float,
):
    base_pred = float(base_pred)
    owners_cost = round(base_pred * (owners_pct / 100.0), 2)
    sst_cost = round(base_pred * (sst_pct / 100.0), 2)
    contingency_cost = round((base_pred + owners_cost) * (cont_pct / 100.0), 2)
    escalation_cost = round((base_pred + owners_cost) * (esc_pct / 100.0), 2)
    eprr_costs = {k: round(base_pred * (float(v) / 100.0), 2) for k, v in (eprr or {}).items()}
    grand_total = round(base_pred + owners_cost + sst_cost + contingency_cost + escalation_cost, 2)
    return owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total


def project_components_df(proj):
    comps = proj.get("components", [])
    rows = []
    for c in comps:
        rows.append(
            {
                "Component": c["component_type"],
                "Dataset": c["dataset"],
                "Model": c.get("model_used", ""),
                "Base CAPEX": float(c["prediction"]),
                "Owner's Cost": float(c["breakdown"]["owners_cost"]),
                "Contingency": float(c["breakdown"]["contingency_cost"]),
                "Escalation": float(c["breakdown"]["escalation_cost"]),
                "SST": float(c["breakdown"]["sst_cost"]),
                "Grand Total": float(c["breakdown"]["grand_total"]),
            }
        )
    return pd.DataFrame(rows)


def project_totals(proj):
    dfc = project_components_df(proj)
    if dfc.empty:
        return {"capex_sum": 0.0, "owners": 0.0, "cont": 0.0, "esc": 0.0, "sst": 0.0, "grand_total": 0.0}
    return {
        "capex_sum": float(dfc["Base CAPEX"].sum()),
        "owners": float(dfc["Owner's Cost"].sum()),
        "cont": float(dfc["Contingency"].sum()),
        "esc": float(dfc["Escalation"].sum()),
        "sst": float(dfc["SST"].sum()),
        "grand_total": float(dfc["Grand Total"].sum()),
    }


# ======================================================================================
# MONTE CARLO HELPERS
# ======================================================================================
def _coerce_float(x, default=np.nan):
    try:
        if x is None:
            return default
        if isinstance(x, str) and x.strip() == "":
            return default
        return float(x)
    except Exception:
        return default


def monte_carlo_component(
    model_pipe: Pipeline,
    feature_cols: list[str],
    base_payload: dict,
    n_sims: int = 5000,
    seed: int = 42,
    feature_sigma_pct: float = 5.0,
    pct_sigma_abs: float = 1.0,
    eprr: dict | None = None,
    sst_pct: float = 0.0,
    owners_pct: float = 0.0,
    cont_pct: float = 0.0,
    esc_pct: float = 0.0,
    normalize_eprr_each_draw: bool = False,
) -> pd.DataFrame:
    rng = np.random.default_rng(int(seed))
    n = int(n_sims)

    base_vec = np.array([_coerce_float(base_payload.get(c), np.nan) for c in feature_cols], dtype=float)
    Xsim = np.tile(base_vec, (n, 1))

    sigma = float(feature_sigma_pct) / 100.0
    if sigma > 0:
        noise = rng.normal(0.0, sigma, size=Xsim.shape)
        mask = ~np.isnan(Xsim)
        Xsim[mask] = Xsim[mask] * (1.0 + noise[mask])

    df_sim = pd.DataFrame(Xsim, columns=feature_cols)
    base_preds = model_pipe.predict(df_sim).astype(float)

    p_sig = float(pct_sigma_abs)
    sst_draw = np.clip(rng.normal(loc=float(sst_pct), scale=p_sig, size=n), 0.0, 100.0)
    own_draw = np.clip(rng.normal(loc=float(owners_pct), scale=p_sig, size=n), 0.0, 100.0)
    con_draw = np.clip(rng.normal(loc=float(cont_pct), scale=p_sig, size=n), 0.0, 100.0)
    esc_draw = np.clip(rng.normal(loc=float(esc_pct), scale=p_sig, size=n), 0.0, 100.0)

    eprr = eprr or {}
    e_keys = list(eprr.keys())
    e_mat = None
    if e_keys:
        e_mat = np.vstack(
            [np.clip(rng.normal(loc=float(eprr.get(k, 0.0)), scale=p_sig, size=n), 0.0, 100.0) for k in e_keys]
        ).T
        if normalize_eprr_each_draw:
            rs = e_mat.sum(axis=1)
            rs[rs == 0] = 1.0
            e_mat = (e_mat / rs[:, None]) * 100.0

    owners_cost = np.round(base_preds * (own_draw / 100.0), 2)
    sst_cost = np.round(base_preds * (sst_draw / 100.0), 2)
    contingency_cost = np.round((base_preds + owners_cost) * (con_draw / 100.0), 2)
    escalation_cost = np.round((base_preds + owners_cost) * (esc_draw / 100.0), 2)
    grand_total = np.round(base_preds + owners_cost + sst_cost + contingency_cost + escalation_cost, 2)

    out = pd.DataFrame(
        {
            "base_pred": base_preds,
            "owners_cost": owners_cost,
            "sst_cost": sst_cost,
            "contingency_cost": contingency_cost,
            "escalation_cost": escalation_cost,
            "grand_total": grand_total,
            "owners_pct": own_draw,
            "sst_pct": sst_draw,
            "cont_pct": con_draw,
            "esc_pct": esc_draw,
        }
    )

    if e_keys:
        for j, k in enumerate(e_keys):
            out[f"eprr_{k}_pct"] = e_mat[:, j]
            out[f"eprr_{k}_cost"] = np.round(base_preds * (e_mat[:, j] / 100.0), 2)

    return out


def scenario_bucket_from_baseline(values: pd.Series, baseline: float, low_cut_pct: float, band_pct: float, high_cut_pct: float):
    baseline = float(baseline) if np.isfinite(baseline) and float(baseline) != 0 else float(values.median())
    pct_delta = (values - baseline) / baseline

    def label(v):
        if v < (-low_cut_pct / 100.0):
            return "Low"
        if (-band_pct / 100.0) <= v <= (band_pct / 100.0):
            return "Base"
        if v > (high_cut_pct / 100.0):
            return "High"
        return "Unbucketed"

    buckets = pct_delta.apply(label)
    return buckets, pct_delta * 100.0


# ---------------------------------------------------------------------------------------
# DATA / MODEL HELPERS
# ---------------------------------------------------------------------------------------
GITHUB_USER = "Hafizuddin-Abd-Rahman-Dev-Upstream"
REPO_NAME = "Cost-Predictor"
BRANCH = "main"
DATA_FOLDER = "pages/data_CAPEX"

MODEL_CANDIDATES = {
    "RandomForest": lambda rs=42: RandomForestRegressor(random_state=rs),
    "GradientBoosting": lambda rs=42: GradientBoostingRegressor(random_state=rs),
    "Ridge": lambda rs=42: Ridge(),
    "Lasso": lambda rs=42: Lasso(),
    "SVR": lambda rs=42: SVR(),
    "DecisionTree": lambda rs=42: DecisionTreeRegressor(random_state=rs),
}

SCALE_MODELS = {"Ridge", "Lasso", "SVR"}  # use scaler


@st.cache_data(ttl=600, show_spinner=False)
def fetch_json(url: str):
    r = requests.get(url, timeout=15)
    r.raise_for_status()
    return r.json()


@st.cache_data(ttl=600, show_spinner=False)
def list_csvs_from_manifest(folder_path: str):
    manifest_url = f"https://raw.githubusercontent.com/{GITHUB_USER}/{REPO_NAME}/{BRANCH}/{folder_path}/files.json"
    try:
        data = fetch_json(manifest_url)
        if isinstance(data, list):
            return [str(x) for x in data]
        return []
    except Exception as e:
        st.error(f"Failed to load CSV manifest: {e}")
        return []


def prepare_numeric_df(df: pd.DataFrame) -> pd.DataFrame:
    num = df.select_dtypes(include=[np.number]).copy()
    num = num.dropna(axis=1, how="all")
    if num.shape[1] < 2:
        raise ValueError("Need at least 2 numeric columns (features + target).")
    return num


def build_X_y(df: pd.DataFrame, target_col: str):
    num = prepare_numeric_df(df)
    if target_col not in num.columns:
        raise ValueError("Selected target is not numeric or not found.")
    X = num.drop(columns=[target_col])
    if X.shape[1] < 1:
        raise ValueError("Need at least 1 numeric feature column.")
    y = num[target_col]
    return X, y


def make_pipeline(model_name: str, random_state=42):
    ctor = MODEL_CANDIDATES[model_name]
    try:
        model = ctor(random_state)
    except TypeError:
        model = ctor()
    steps = [("imputer", SimpleImputer(strategy="median"))]
    if model_name in SCALE_MODELS:
        steps.append(("scaler", MinMaxScaler()))
    steps.append(("model", model))
    return Pipeline(steps)


def evaluate_models(X, y, test_size=0.2, random_state=42):
    Xtr, Xte, ytr, yte = train_test_split(X, y, test_size=test_size, random_state=random_state)

    rows = []
    best_name = None
    best_r2 = -np.inf
    best_rmse = None

    for name in MODEL_CANDIDATES.keys():
        pipe = make_pipeline(name, random_state=random_state)
        pipe.fit(Xtr, ytr)
        yhat = pipe.predict(Xte)
        rmse = float(np.sqrt(mean_squared_error(yte, yhat)))
        r2 = float(r2_score(yte, yhat))
        rows.append({"model": name, "rmse": rmse, "r2": r2})
        if r2 > best_r2:
            best_r2 = r2
            best_rmse = rmse
            best_name = name

    rows_sorted = sorted(rows, key=lambda d: d["r2"], reverse=True)
    metrics = {"best_model": best_name, "rmse": best_rmse, "r2": best_r2, "models": rows_sorted}
    return metrics


@st.cache_resource(show_spinner=False)
def train_best_model_cached(df: pd.DataFrame, target_col: str, test_size: float, random_state: int, dataset_key: str):
    X, y = build_X_y(df, target_col)
    metrics = evaluate_models(X, y, test_size=test_size, random_state=random_state)
    best_name = metrics.get("best_model") or "RandomForest"
    best_pipe = make_pipeline(best_name, random_state=random_state)
    best_pipe.fit(X, y)
    return best_pipe, metrics, list(X.columns), y.name, best_name


def single_prediction(model_pipe: Pipeline, feature_cols: list[str], payload: dict):
    row = {}
    for c in feature_cols:
        v = payload.get(c, np.nan)
        try:
            if v is None or (isinstance(v, str) and v.strip() == ""):
                row[c] = np.nan
            else:
                row[c] = float(v)
        except Exception:
            row[c] = np.nan
    df_in = pd.DataFrame([row], columns=feature_cols)
    return float(model_pipe.predict(df_in)[0])


@st.cache_data(show_spinner=False, ttl=600)
def knn_impute_numeric(df: pd.DataFrame, k: int = 5):
    num = prepare_numeric_df(df)
    arr = KNNImputer(n_neighbors=k).fit_transform(num)
    return pd.DataFrame(arr, columns=num.columns)


# =======================================================================================
# EXPORT HELPERS (Excel / PPT)
# =======================================================================================
def _format_ws_money(ws, start_row=2):
    ws.freeze_panes = "A2"
    for c in range(1, ws.max_column + 1):
        ws.column_dimensions[get_column_letter(c)].width = 18
    for r in range(start_row, ws.max_row + 1):
        for c in range(1, ws.max_column + 1):
            cell = ws.cell(r, c)
            if isinstance(cell.value, (int, float)) and c >= 3:
                cell.number_format = "#,##0.00"


def create_project_excel_report_capex(project_name, proj, currency=""):
    output = io.BytesIO()
    comps_df = project_components_df(proj)

    if comps_df.empty:
        with pd.ExcelWriter(output, engine="openpyxl") as writer:
            pd.DataFrame({"Info": [f"No components for project {project_name}"]}).to_excel(
                writer, sheet_name="Summary", index=False
            )
        output.seek(0)
        return output

    totals = project_totals(proj)

    summary_df = comps_df.copy()
    summary_df.loc[len(summary_df)] = {
        "Component": "TOTAL",
        "Dataset": "",
        "Model": "",
        "Base CAPEX": totals["capex_sum"],
        "Owner's Cost": totals["owners"],
        "Contingency": totals["cont"],
        "Escalation": totals["esc"],
        "SST": totals["sst"],
        "Grand Total": totals["grand_total"],
    }

    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df.to_excel(writer, sheet_name="Summary", index=False)
        ws = writer.sheets["Summary"]
        _format_ws_money(ws)

        max_row = ws.max_row
        max_col = ws.max_column

        for col_idx in range(4, max_col + 1):
            col_letter = get_column_letter(col_idx)
            ws.conditional_formatting.add(
                f"{col_letter}2:{col_letter}{max_row-1}",
                ColorScaleRule(
                    start_type="percentile",
                    start_value=10,
                    start_color="FFE0F7FA",
                    mid_type="percentile",
                    mid_value=50,
                    mid_color="FF80DEEA",
                    end_type="percentile",
                    end_value=90,
                    end_color="FF00838F",
                ),
            )

        chart = BarChart()
        chart.title = "Grand Total by Component"
        data = Reference(ws, min_col=10, max_col=10, min_row=1, max_row=max_row - 1)
        cats = Reference(ws, min_col=1, min_row=2, max_row=max_row - 1)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.y_axis.title = f"Cost ({currency})".strip()
        chart.x_axis.title = "Component"
        chart.height = 10
        chart.width = 18
        ws.add_chart(chart, "L2")

    output.seek(0)
    return output


def create_project_pptx_report_capex(project_name, proj, currency=""):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    layout_title_only = prs.slide_layouts[5]
    layout_title_content = prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout_title_only)
    title = slide.shapes.title
    title.text = f"CAPEX Project Report\n{project_name}"
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 161, 155)

    comps_df = project_components_df(proj)
    comps = proj.get("components", [])
    totals = project_totals(proj)

    slide = prs.slides.add_slide(layout_title_content)
    slide.shapes.title.text = "Executive Summary"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    lines = [
        f"Project: {project_name}",
        f"Total Components: {len(comps)}",
        f"Total Base CAPEX: {currency} {totals['capex_sum']:,.2f}",
        f"Total Grand Total (incl. SST): {currency} {totals['grand_total']:,.2f}",
    ]
    for i, line in enumerate(lines):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.text = line
        para.font.size = Pt(16)

    if not comps_df.empty:
        fig, ax = plt.subplots(figsize=(7, 4))
        ax.bar(comps_df["Component"], comps_df["Grand Total"])
        ax.set_title("Grand Total by Component")
        ax.set_ylabel(f"Cost ({currency})".strip())
        ax.tick_params(axis="x", rotation=25)
        ax.grid(axis="y", linestyle="--", alpha=0.4)
        fig.tight_layout()

        img_stream = io.BytesIO()
        fig.savefig(img_stream, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        img_stream.seek(0)

        slide2 = prs.slides.add_slide(layout_title_only)
        slide2.shapes.title.text = "Grand Total by Component"
        slide2.shapes.add_picture(img_stream, Inches(0.7), Inches(1.5), width=Inches(8.6))

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ---------------------------------------------------------------------------------------
# TOP-LEVEL TABS (MC is its OWN TAB)
# ---------------------------------------------------------------------------------------
tab_data, tab_pb, tab_mc, tab_compare = st.tabs(
    ["📊 Data", "🏗️ Project Builder", "🎲 Monte Carlo", "🔀 Compare Projects"]
)

# =======================================================================================
# DATA TAB
# =======================================================================================
with tab_data:
    st.markdown("### 📁 Data")

    c1, c2 = st.columns([1.2, 1])
    with c1:
        data_source = st.radio("Choose data source", ["Upload CSV", "Load from Server"], horizontal=True, key="data_source")
    with c2:
        st.caption("Tip: Streamlit Cloud needs requirements.txt for sklearn.")

    uploaded_files = []
    if data_source == "Upload CSV":
        uploaded_files = st.file_uploader(
            "Upload CSV files",
            type="csv",
            accept_multiple_files=True,
            key=f"csv_uploader_{st.session_state.uploader_nonce}",
        )
    else:
        github_csvs = list_csvs_from_manifest(DATA_FOLDER)
        if github_csvs:
            selected_file = st.selectbox("Choose CSV from GitHub", github_csvs, key="github_csv_select")
            if st.button("Load selected CSV", key="btn_load_github_csv"):
                raw_url = f"https://raw.githubusercontent.com/{GITHUB_USER}/{REPO_NAME}/{BRANCH}/{DATA_FOLDER}/{selected_file}"
                try:
                    df = pd.read_csv(raw_url)
                    st.session_state.datasets[selected_file] = df
                    st.session_state.predictions.setdefault(selected_file, [])
                    toast(f"Loaded from GitHub: {selected_file}")
                    st.rerun()
                except Exception as e:
                    st.error(f"Error loading CSV: {e}")
        else:
            st.info("No CSV files found in GitHub folder (files.json missing or empty).")

    if uploaded_files:
        for up in uploaded_files:
            if up.name not in st.session_state.datasets:
                try:
                    df = pd.read_csv(up)
                    st.session_state.datasets[up.name] = df
                    st.session_state.predictions.setdefault(up.name, [])
                except Exception as e:
                    st.error(f"Failed to read {up.name}: {e}")
        toast("Dataset(s) added.")

    st.divider()

    if not st.session_state.datasets:
        st.info("Upload or load a dataset to proceed.")
        st.stop()

    ds_name_data = st.selectbox("Active dataset", list(st.session_state.datasets.keys()), key="active_dataset")
    df_active = st.session_state.datasets[ds_name_data]

    num_cols = df_active.select_dtypes(include=[np.number]).columns.tolist()
    if len(num_cols) < 2:
        st.warning("This dataset has < 2 numeric columns. Model requires numeric features + numeric target.")
        st.stop()

    target_key = f"target_col__{ds_name_data}"
    if target_key not in st.session_state:
        st.session_state[target_key] = num_cols[-1]

    target_col_active = st.selectbox(
        "Target (Cost) column",
        options=num_cols,
        index=num_cols.index(st.session_state[target_key]) if st.session_state[target_key] in num_cols else len(num_cols) - 1,
        key=f"{target_key}__selector",
    )
    # store back
    st.session_state[target_key] = target_col_active

    currency_active = get_currency_symbol(df_active, target_col_active)

    a, b, c = st.columns(3)
    a.metric("Rows", f"{df_active.shape[0]:,}")
    b.metric("Columns", f"{df_active.shape[1]:,}")
    c.metric("Currency", currency_active or "—")

    with st.expander("Preview (first 15 rows)", expanded=False):
        st.dataframe(df_active.head(15), use_container_width=True)

    st.divider()

    st.markdown("### ⚙️ Model Training")
    test_size = st.slider("Test size", 0.1, 0.5, 0.2, 0.05, key="train_test_size")
    if st.button("Run training", key="btn_run_training"):
        try:
            with st.spinner("Training model (cached)..."):
                pipe, metrics, feat_cols, y_name, best_name = train_best_model_cached(
                    df_active,
                    target_col_active,
                    test_size=float(test_size),
                    random_state=42,
                    dataset_key=ds_name_data,
                )
            st.session_state._last_metrics = metrics
            toast("Training complete.")
            m1, m2, m3 = st.columns(3)
            m1.metric("RMSE (best)", f"{metrics['rmse']:,.2f}")
            m2.metric("R² (best)", f"{metrics['r2']:.3f}")
            m3.metric("Best Model", best_name)
        except Exception as e:
            st.error(f"Training failed: {e}")

# =======================================================================================
# PROJECT BUILDER TAB
# =======================================================================================
with tab_pb:
    st.markdown("### 🏗️ Project Builder")

    if not st.session_state.datasets:
        st.info("No dataset. Go to **Data** tab to upload or load.")
        st.stop()

    colA, colB = st.columns([2, 1])
    with colA:
        new_project_name = st.text_input("New Project Name", placeholder="e.g., CAPEX 2026", key="pb_new_project_name")
    with colB:
        if new_project_name and new_project_name not in st.session_state.projects:
            if st.button("Create Project", key="pb_create_project_btn"):
                st.session_state.projects[new_project_name] = {"components": [], "totals": {}, "currency": ""}
                toast(f"Project '{new_project_name}' created.")
                st.rerun()

    if not st.session_state.projects:
        st.info("Create a project above, then add components.")
        st.stop()

    proj_sel = st.selectbox("Select project to work on", list(st.session_state.projects.keys()), key="pb_project_select")

    ds_names = sorted(st.session_state.datasets.keys())
    dataset_for_comp = st.selectbox("Dataset for this component", ds_names, key=f"pb_dataset_for_component__{proj_sel}")
    df_comp = st.session_state.datasets[dataset_for_comp]

    target_col_comp = st.session_state.get(f"target_col__{dataset_for_comp}", None)
    if not target_col_comp:
        # fallback
        num_cols_comp = df_comp.select_dtypes(include=[np.number]).columns.tolist()
        target_col_comp = num_cols_comp[-1] if num_cols_comp else None

    curr_ds = get_currency_symbol(df_comp, target_col_comp)

    default_label = st.session_state.component_labels.get(dataset_for_comp, "")
    component_type = st.text_input(
        "Component type (Asset / Scope)",
        value=(default_label or "Platform / Pipeline / Subsea / Well"),
        key=f"pb_component_type__{proj_sel}__{dataset_for_comp}",
    )

    try:
        pipe_c, _, feat_cols_c, y_name_c, best_name_c = train_best_model_cached(
            df_comp,
            target_col_comp,
            test_size=0.2,
            random_state=42,
            dataset_key=dataset_for_comp,
        )
    except Exception as e:
        st.error(f"Component model setup failed: {e}")
        st.stop()

    st.markdown("**Component Feature Inputs (1 row)**")

    # ✅ KEY FIX for StreamlitDuplicateElementId:
    # Unique key includes project + dataset + a nonce that changes when needed.
    comp_input_key = f"pb_input_row__{proj_sel}__{dataset_for_comp}"
    if comp_input_key not in st.session_state:
        st.session_state[comp_input_key] = {c: np.nan for c in feat_cols_c}

    comp_row_df = pd.DataFrame([st.session_state[comp_input_key]], columns=feat_cols_c)
    comp_edited = st.data_editor(
        comp_row_df,
        num_rows="fixed",
        use_container_width=True,
        key=f"pb_data_editor__{proj_sel}__{dataset_for_comp}__{st.session_state._widget_nonce}",
    )
    comp_payload = comp_edited.iloc[0].to_dict()

    st.markdown("---")
    st.markdown("**Cost Percentage Inputs**")
    cp1, cp2 = st.columns(2)
    with cp1:
        st.markdown("EPRR (%)")
        eng_pb = st.number_input("Engineering", 0.0, 100.0, 12.0, 1.0, key=f"pb_eng__{proj_sel}")
        prep_pb = st.number_input("Preparation", 0.0, 100.0, 7.0, 1.0, key=f"pb_prep__{proj_sel}")
        remv_pb = st.number_input("Removal", 0.0, 100.0, 54.0, 1.0, key=f"pb_remv__{proj_sel}")
        remd_pb = st.number_input("Remediation", 0.0, 100.0, 27.0, 1.0, key=f"pb_remd__{proj_sel}")

        eprr_pb = {"Engineering": eng_pb, "Preparation": prep_pb, "Removal": remv_pb, "Remediation": remd_pb}
        eprr_total_pb = sum(eprr_pb.values())
        st.caption(f"EPRR total: **{eprr_total_pb:.2f}%**")
        apply_norm_pb = st.checkbox("Normalize EPRR to 100%", value=False, key=f"pb_norm__{proj_sel}")
        if apply_norm_pb and eprr_total_pb > 0 and abs(eprr_total_pb - 100.0) > 1e-6:
            eprr_pb, _ = normalize_to_100(eprr_pb)

    with cp2:
        st.markdown("Financial (%)")
        sst_pb = st.number_input("SST", 0.0, 100.0, 0.0, 0.5, key=f"pb_sst__{proj_sel}")
        owners_pb = st.number_input("Owner's Cost", 0.0, 100.0, 0.0, 0.5, key=f"pb_owners__{proj_sel}")
        cont_pb = st.number_input("Contingency", 0.0, 100.0, 0.0, 0.5, key=f"pb_cont__{proj_sel}")
        esc_pb = st.number_input("Escalation & Inflation", 0.0, 100.0, 0.0, 0.5, key=f"pb_esc__{proj_sel}")

    if st.button("➕ Predict & Add Component", key=f"pb_add_comp__{proj_sel}__{dataset_for_comp}"):
        try:
            base_pred = single_prediction(pipe_c, feat_cols_c, comp_payload)
            owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total = cost_breakdown(
                base_pred, eprr_pb, sst_pb, owners_pb, cont_pb, esc_pb
            )

            comp_entry = {
                "component_type": component_type or default_label or "Component",
                "dataset": dataset_for_comp,
                "model_used": best_name_c,
                "inputs": {k: comp_payload.get(k, np.nan) for k in feat_cols_c},
                "feature_cols": list(feat_cols_c),  # for MC reuse
                "prediction": float(base_pred),
                "breakdown": {
                    "eprr_costs": eprr_costs,
                    "eprr_pct": eprr_pb,
                    "sst_cost": float(sst_cost),
                    "owners_cost": float(owners_cost),
                    "contingency_cost": float(contingency_cost),
                    "escalation_cost": float(escalation_cost),
                    "grand_total": float(grand_total),
                    "target_col": y_name_c,
                    "sst_pct": float(sst_pb),
                    "owners_pct": float(owners_pb),
                    "cont_pct": float(cont_pb),
                    "esc_pct": float(esc_pb),
                },
            }

            st.session_state.projects[proj_sel]["components"].append(comp_entry)
            st.session_state.component_labels[dataset_for_comp] = component_type or default_label
            st.session_state.projects[proj_sel]["currency"] = curr_ds

            # bump nonce to avoid any lingering widget-id conflicts after rerun
            st.session_state._widget_nonce += 1

            toast(f"Component added to project '{proj_sel}'.")
            st.rerun()
        except Exception as e:
            st.error(f"Failed to add component: {e}")

    st.markdown("---")
    st.markdown("### Current Project Overview")

    proj = st.session_state.projects[proj_sel]
    comps = proj.get("components", [])
    if not comps:
        st.info("No components yet. Add at least one above.")
        st.stop()

    dfc = project_components_df(proj)
    curr = proj.get("currency", "") or curr_ds

    st.dataframe(
        dfc.style.format(
            {
                "Base CAPEX": "{:,.2f}",
                "Owner's Cost": "{:,.2f}",
                "Contingency": "{:,.2f}",
                "Escalation": "{:,.2f}",
                "SST": "{:,.2f}",
                "Grand Total": "{:,.2f}",
            }
        ),
        use_container_width=True,
        height=320,
    )

    t = project_totals(proj)
    proj["totals"] = {"capex_sum": t["capex_sum"], "grand_total": t["grand_total"]}

    m1, m2, m3 = st.columns(3)
    m1.metric("Project CAPEX (Base)", f"{curr} {t['capex_sum']:,.2f}")
    m2.metric("Project SST", f"{curr} {t['sst']:,.2f}")
    m3.metric("Project Grand Total (incl. SST)", f"{curr} {t['grand_total']:,.2f}")

    st.markdown("#### Components list")
    for idx, c in enumerate(list(comps)):
        a, b, d = st.columns([4, 2, 1])
        a.write(f"**{c['component_type']}** — *{c['dataset']}* — {c.get('model_used','')}")
        b.write(f"GT: {curr} {c['breakdown']['grand_total']:,.2f}")
        if d.button("🗑️ Remove", key=f"pb_del_comp__{proj_sel}__{idx}"):
            comps.pop(idx)
            st.session_state._widget_nonce += 1
            toast("Component removed.", "🗑️")
            st.rerun()

    st.markdown("---")
    st.markdown("#### Export Project")
    col1, col2, col3 = st.columns(3)

    with col1:
        excel_report = create_project_excel_report_capex(proj_sel, proj, curr)
        st.download_button(
            "⬇️ Download Project Excel",
            data=excel_report,
            file_name=f"{proj_sel}_CAPEX_Report.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key=f"dl_proj_excel__{proj_sel}",
        )

    with col2:
        pptx_report = create_project_pptx_report_capex(proj_sel, proj, curr)
        st.download_button(
            "⬇️ Download Project PowerPoint",
            data=pptx_report,
            file_name=f"{proj_sel}_CAPEX_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"dl_proj_pptx__{proj_sel}",
        )

    with col3:
        st.download_button(
            "⬇️ Download Project (JSON)",
            data=json.dumps(proj, indent=2, default=float),
            file_name=f"{proj_sel}.json",
            mime="application/json",
            key=f"dl_proj_json__{proj_sel}",
        )


# =======================================================================================
# 🎲 MONTE CARLO TAB
# =======================================================================================
with tab_mc:
    st.markdown("### 🎲 Monte Carlo")
    st.caption("Simulate uncertainty per component and roll-up to project Grand Total distribution.")

    if not st.session_state.projects:
        st.info("No projects found. Create a project first in **🏗️ Project Builder**.")
        st.stop()

    proj_sel_mc = st.selectbox("Select project", list(st.session_state.projects.keys()), key="mc_project_select")

    proj = st.session_state.projects[proj_sel_mc]
    comps = proj.get("components", [])
    if not comps:
        st.warning("This project has no components. Add components in **🏗️ Project Builder** first.")
        st.stop()

    t_mc = project_totals(proj)
    curr = proj.get("currency", "") or ""
    baseline_gt = float(t_mc["grand_total"])
    st.info(f"Baseline (current Project Grand Total): **{curr} {baseline_gt:,.2f}**")

    st.markdown("#### Simulation settings")
    mcA, mcB, mcC = st.columns(3)
    with mcA:
        mc_n_sims = st.number_input("Simulations", 500, 50000, 5000, 500, key=f"mc_n__{proj_sel_mc}")
        mc_seed = st.number_input("Random seed", 0, 999999, 42, 1, key=f"mc_seed__{proj_sel_mc}")
    with mcB:
        mc_feat_sigma = st.slider("Feature uncertainty (±% stdev)", 0.0, 30.0, 5.0, 0.5, key=f"mc_feat__{proj_sel_mc}")
        mc_pct_sigma = st.slider("Percent uncertainty (± abs stdev)", 0.0, 10.0, 1.0, 0.1, key=f"mc_pct__{proj_sel_mc}")
    with mcC:
        mc_norm_eprr = st.checkbox("Normalize EPRR to 100% each simulation", False, key=f"mc_norm__{proj_sel_mc}")
        mc_budget = st.number_input(
            "Budget threshold (Project Grand Total)",
            min_value=0.0,
            value=float(baseline_gt),
            step=1000.0,
            key=f"mc_budget__{proj_sel_mc}",
        )

    st.markdown("#### Scenario buckets vs baseline")
    sb1, sb2, sb3 = st.columns(3)
    with sb1:
        mc_low = st.slider("Low < baseline by (%)", 0, 50, 10, 1, key=f"mc_low__{proj_sel_mc}")
    with sb2:
        mc_band = st.slider("Base band ± (%)", 1, 50, 10, 1, key=f"mc_band__{proj_sel_mc}")
    with sb3:
        mc_high = st.slider("High > baseline by (%)", 0, 50, 10, 1, key=f"mc_high__{proj_sel_mc}")

    if st.button("Run Monte Carlo", type="primary", key=f"mc_run__{proj_sel_mc}"):
        try:
            with st.spinner("Running Monte Carlo..."):
                n = int(mc_n_sims)
                project_gt = np.zeros(n, dtype=float)
                comp_summ_rows = []

                for idx, comp in enumerate(comps):
                    ds_name = comp["dataset"]
                    df_ds = st.session_state.datasets.get(ds_name)
                    if df_ds is None:
                        raise ValueError(f"Dataset not found in session: {ds_name}")

                    target_col = comp["breakdown"].get("target_col")
                    if not target_col:
                        raise ValueError(f"Component '{comp['component_type']}' missing breakdown.target_col")

                    pipe_tmp, _, feat_cols_tmp, _, _ = train_best_model_cached(
                        df_ds, target_col, test_size=0.2, random_state=42, dataset_key=ds_name
                    )

                    feat_cols = comp.get("feature_cols") or feat_cols_tmp
                    payload = comp.get("inputs") or {}

                    eprr_pct = comp["breakdown"].get("eprr_pct", {})
                    sst_pct_c = float(comp["breakdown"].get("sst_pct", 0.0))
                    owners_pct_c = float(comp["breakdown"].get("owners_pct", 0.0))
                    cont_pct_c = float(comp["breakdown"].get("cont_pct", 0.0))
                    esc_pct_c = float(comp["breakdown"].get("esc_pct", 0.0))

                    comp_seed = int(mc_seed) + (idx + 1) * 101

                    df_mc_c = monte_carlo_component(
                        model_pipe=pipe_tmp,
                        feature_cols=list(feat_cols),
                        base_payload=payload,
                        n_sims=n,
                        seed=comp_seed,
                        feature_sigma_pct=float(mc_feat_sigma),
                        pct_sigma_abs=float(mc_pct_sigma),
                        eprr=eprr_pct,
                        sst_pct=sst_pct_c,
                        owners_pct=owners_pct_c,
                        cont_pct=cont_pct_c,
                        esc_pct=esc_pct_c,
                        normalize_eprr_each_draw=bool(mc_norm_eprr),
                    )

                    project_gt += df_mc_c["grand_total"].to_numpy(dtype=float)

                    comp_summ_rows.append(
                        {
                            "Component": comp["component_type"],
                            "Dataset": ds_name,
                            "P50": float(df_mc_c["grand_total"].quantile(0.50)),
                            "P80": float(df_mc_c["grand_total"].quantile(0.80)),
                            "P90": float(df_mc_c["grand_total"].quantile(0.90)),
                        }
                    )

                df_proj_mc = pd.DataFrame({"project_grand_total": project_gt})
                buckets, pct_delta = scenario_bucket_from_baseline(
                    df_proj_mc["project_grand_total"], baseline_gt, mc_low, mc_band, mc_high
                )
                df_proj_mc["Scenario"] = buckets
                df_proj_mc["%Δ vs baseline"] = pct_delta

                p50 = float(df_proj_mc["project_grand_total"].quantile(0.50))
                p80 = float(df_proj_mc["project_grand_total"].quantile(0.80))
                p90 = float(df_proj_mc["project_grand_total"].quantile(0.90))
                exceed_prob = float((df_proj_mc["project_grand_total"] > float(mc_budget)).mean()) * 100.0

                df_comp_mc = pd.DataFrame(comp_summ_rows)

            st.markdown("#### Summary")
            m1, m2, m3, m4 = st.columns(4)
            m1.metric("P50 Project GT", f"{curr} {p50:,.2f}")
            m2.metric("P80 Project GT", f"{curr} {p80:,.2f}")
            m3.metric("P90 Project GT", f"{curr} {p90:,.2f}")
            m4.metric("P(> Budget)", f"{exceed_prob:.1f}%")

            st.plotly_chart(
                px.histogram(df_proj_mc, x="project_grand_total", nbins=60, title="Project Grand Total distribution"),
                use_container_width=True,
            )

            bucket_counts = df_proj_mc["Scenario"].value_counts().reset_index()
            bucket_counts.columns = ["Scenario", "Count"]
            st.plotly_chart(px.bar(bucket_counts, x="Scenario", y="Count", title="Scenario bucket counts"), use_container_width=True)

            st.markdown("#### Component summary (P50/P80/P90)")
            st.dataframe(
                df_comp_mc.style.format({"P50": "{:,.2f}", "P80": "{:,.2f}", "P90": "{:,.2f}"}),
                use_container_width=True,
            )

            with st.expander("Show Monte Carlo table (first 200 rows)", expanded=False):
                st.dataframe(df_proj_mc.head(200), use_container_width=True)

            st.markdown("#### Download MC results")
            csv_proj = df_proj_mc.to_csv(index=False).encode("utf-8")
            csv_comp = df_comp_mc.to_csv(index=False).encode("utf-8")

            c1, c2, c3 = st.columns(3)
            c1.download_button(
                "⬇️ Project MC (CSV)",
                data=csv_proj,
                file_name=f"{proj_sel_mc}_mc_project.csv",
                mime="text/csv",
                key=f"dl_mc_proj_csv__{proj_sel_mc}",
            )
            c2.download_button(
                "⬇️ Component MC Summary (CSV)",
                data=csv_comp,
                file_name=f"{proj_sel_mc}_mc_components.csv",
                mime="text/csv",
                key=f"dl_mc_comp_csv__{proj_sel_mc}",
            )

            bio_xlsx = io.BytesIO()
            with pd.ExcelWriter(bio_xlsx, engine="openpyxl") as writer:
                df_proj_mc.to_excel(writer, sheet_name="Project_MC", index=False)
                df_comp_mc.to_excel(writer, sheet_name="Component_Summary", index=False)
            bio_xlsx.seek(0)

            zip_bio = io.BytesIO()
            with zipfile.ZipFile(zip_bio, "w", zipfile.ZIP_DEFLATED) as zf:
                zf.writestr(f"{proj_sel_mc}_mc_project.csv", csv_proj)
                zf.writestr(f"{proj_sel_mc}_mc_components.csv", csv_comp)
                zf.writestr(f"{proj_sel_mc}_mc_results.xlsx", bio_xlsx.getvalue())
            zip_bio.seek(0)

            c3.download_button(
                "⬇️ MC Results (ZIP)",
                data=zip_bio.getvalue(),
                file_name=f"{proj_sel_mc}_mc_results.zip",
                mime="application/zip",
                key=f"dl_mc_zip__{proj_sel_mc}",
            )

        except Exception as e:
            st.error(f"Monte Carlo failed: {e}")


# =======================================================================================
# COMPARE PROJECTS TAB
# =======================================================================================
with tab_compare:
    st.markdown("### 🔀 Compare Projects")

    proj_names = list(st.session_state.projects.keys())
    if len(proj_names) < 2:
        st.info("Create at least two projects in the Project Builder tab to compare.")
        st.stop()

    compare_sel = st.multiselect("Select projects", proj_names, default=proj_names[:2], key="compare_sel")
    if len(compare_sel) < 2:
        st.warning("Select at least two projects.")
        st.stop()

    rows = []
    for p in compare_sel:
        proj = st.session_state.projects[p]
        t = project_totals(proj)
        rows.append(
            {
                "Project": p,
                "Components": len(proj.get("components", [])),
                "CAPEX Sum": t["capex_sum"],
                "Owner": t["owners"],
                "Contingency": t["cont"],
                "Escalation": t["esc"],
                "SST": t["sst"],
                "Grand Total": t["grand_total"],
                "Currency": proj.get("currency", ""),
            }
        )

    df_proj = pd.DataFrame(rows)
    st.dataframe(
        df_proj[["Project", "Components", "CAPEX Sum", "SST", "Grand Total", "Currency"]].style.format(
            {"CAPEX Sum": "{:,.2f}", "SST": "{:,.2f}", "Grand Total": "{:,.2f}"}
        ),
        use_container_width=True,
    )

    st.plotly_chart(px.bar(df_proj, x="Project", y="Grand Total", title="Grand Total by Project"), use_container_width=True)
    df_melt = df_proj.melt(
        id_vars=["Project"],
        value_vars=["CAPEX Sum", "Owner", "Contingency", "Escalation", "SST"],
        var_name="Cost Type",
        value_name="Value",
    )
    st.plotly_chart(px.bar(df_melt, x="Project", y="Value", color="Cost Type", barmode="stack", title="Cost composition"),
                    use_container_width=True)
